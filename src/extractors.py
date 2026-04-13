
"""
extractors.py
=============
阶段二：各文件类型的文本提取函数。

架构分为两层：
  1. 提取层（返回 section 列表）：parse_*() / _extract_*_sections()
  2. 分块层（返回 chunk 列表）：extract_*() — 内部调用提取层 + make_chunks_from_sections

公开 API：
  extract_sections(file_record)  → List[dict] | None   # 阶段 2a：纯提取
  extract_*(file_record)         → List[dict]           # 阶段 2a+2b：提取 + 分块

Section 字段：
    section_id, page_or_sheet, text, section_title

ChunkRecord 字段：
    chunk_id, parent_id, file_path, file_name, folder_code,
    page_or_sheet, chunk_index, text, parent_text, char_count

已实现：
    extract_pdf(file_record)   → List[dict]  # 步骤 2-1（含智谱线上 OCR）
    extract_docx(file_record)  → List[dict]  # 步骤 2-2
    extract_doc(file_record)   → List[dict]  # 步骤 2-3（LibreOffice 转换）
    extract_xlsx(file_record)  → List[dict]  # 步骤 2-4
    extract_xls(file_record)   → List[dict]  # 步骤 2-4
    extract_pptx(file_record)  → List[dict]  # 步骤 2-5
    extract_ppt(file_record)   → List[dict]  # 步骤 2-5（LibreOffice 转换）
    extract_image(file_record) → List[dict]  # 步骤 2-6（JPG/PNG，智谱 OCR）
"""

import os
import re
import threading
from typing import List

from config import (
    # 分块参数（全局默认 + 各类型覆盖）
    chunk_params,
    # PDF 分类阈值
    PDF_SAMPLE_PAGES,
    PDF_PPT_ASPECT_RATIO,
    PDF_PPT_AVG_BLOCK_CHARS,
    PDF_PPT_AVG_IMAGES,
    PDF_TITLE_SIZE_PERCENTILE,
    PDF_TITLE_MAX_CHARS,
    PDF_TITLE_RATIO_MAX,
    PDF_TITLE_MIN_COUNT,
    # PDF 提取参数
    PDF_PAGE_MIN_CHARS,
    TITLE_REBUILD_MODE,
    # DOCX 标题识别
    DOCX_HEADING_MAX_CHARS,
    # LibreOffice 超时
    SOFFICE_DOC_TIMEOUT,
    SOFFICE_PPT_TIMEOUT,
    # VLM 图片分类与描述
    VLM_MODEL,
    # VLM 图片过滤与处理
    IMAGE_MIN_WIDTH,
    IMAGE_MIN_HEIGHT,
    IMAGE_MIN_BYTES,
    VLM_MAX_IMAGE_PX,
    # VLM 缓存
    VLM_CACHE_PATH,
    # LLM API（标题层级 LLM 重建用）
    DRAFT_LLM_BASE_URL,
    DRAFT_LLM_API_KEY,
    DRAFT_LLM_MODEL,
    DRAFT_ENABLE_THINKING,
    # 表格处理优化参数（Phase 1）
    TABLE_MAX_ROWS,
    SHORT_TEXT_THRESHOLD,
)


# ==============================================================================
# 2-0  基础设施：字符统计 + 分块
# ==============================================================================

def count_meaningful_chars(text: str) -> int:
    """只统计中文、英文字母、数字，过滤乱码字符（如私有区字符）。"""
    return len(re.findall(r'[\u4e00-\u9fff\u3400-\u4dbfa-zA-Z0-9]', text))


# 分块时按优先级依次尝试的分隔符列表；空字符串 '' 代表按单字符切
SPLIT_SEPARATORS = ['\n\n', '\n', '。', '，', '']

# 分块参数统一在 src/config.py 中管理（CHUNK_MAX_SIZE / CHUNK_MIN_SIZE 及各类型覆盖值）
# 各 parse_* 函数通过 chunk_params('<type>') 获取对应的 (max_size, min_size)


# ==============================================================================
# 2-0b  HTML 表格 → Markdown 转换
# ==============================================================================

def html_table_to_markdown(html_text: str) -> tuple:
    """
    将文本中的 HTML 表格转换为 Markdown 格式，同时保留原始 HTML。

    - 保留非表格部分的文本
    - 多个表格依次转换
    - 正确识别表头（第一行作为列名）
    - 使用 pandas.read_html() + DataFrame.to_markdown() 组合

    技术背景：
        SDK 路径（PP-DocLayout-V3）识别表格后用 "Table Recognition:" prompt
        输出 HTML <table> 格式。HTML 标签占比约 55-60%，信噪比低（~0.8）。
        转换为 Markdown 后标签占比降至 15-20%，信噪比提升至 ~4.0。

    表头处理：
        GLM-OCR 输出的表格使用 <td> 而非 <th>，pandas 默认把第一行当数据。
        本函数会将第一行数据提升为列名（header=0），正确还原表格结构。

    返回：
        (converted_text, table_html_list)
        - converted_text: 转换后的文本（HTML 表格已替换为 Markdown）
        - table_html_list: 原始 HTML 表格列表（用于 Web 渲染）

    依赖：
        pip install tabulate  # pandas.to_markdown() 依赖
    """
    if '<table' not in html_text.lower():
        return html_text, []

    try:
        import pandas as pd
        from io import StringIO
    except ImportError:
        # pandas 不可用时返回原文
        return html_text, []

    result = html_text
    table_html_list = []

    # 处理完整的表格（<table>...</table>）
    complete_pattern = re.compile(r'<table[^>]*>.*?</table>', re.DOTALL | re.IGNORECASE)

    for match in complete_pattern.finditer(html_text):
        table_html = match.group()
        table_html_list.append(table_html)  # 保留原始 HTML
        md = _convert_single_table(table_html, pd, StringIO)
        if md:
            result = result.replace(table_html, '\n' + md + '\n')

    return result, table_html_list


def _convert_single_table(table_html: str, pd, StringIO) -> str:
    """
    将单个 HTML 表格转换为 Markdown。

    参数:
        table_html: 完整的 <table>...</table> HTML 字符串
        pd: pandas 模块引用
        StringIO: io.StringIO 类引用

    返回:
        Markdown 格式的表格字符串，失败时返回 None
    """
    try:
        # pandas.read_html() 需要 StringIO 包装
        # header=0 表示第一行作为列名（表头）
        dfs = pd.read_html(StringIO(table_html), header=0)
        if not dfs:
            return None

        df = dfs[0].fillna('')

        # 将所有列转为字符串，避免数字格式问题
        df = df.astype(str)
        # 将 'nan' 替换为空字符串
        df = df.replace('nan', '')

        # 尝试转换为 Markdown
        try:
            md = df.to_markdown(index=False)
            return md
        except ImportError:
            # tabulate 未安装
            return None

    except Exception as e:
        print(f"[WARN] 表格转 Markdown 失败：{e}")
        return None


def recursive_split(text: str, max_size: int, min_size: int) -> List[str]:
    """
    递归分块：按 SPLIT_SEPARATORS 优先级将 text 切成 ≤ max_size 的片段，
    然后贪心合并，使每个 chunk 尽量接近 max_size。

    合并策略：
      - 贪心：只要"当前累积 + 下一片段"≤ max_size，就持续合并
      - 兜底：合并后仍 < min_size 的孤立片段与前一个合并（即使超过 max_size 也接受，
              但 _split 保证单片段本身不超 max_size，所以不会出现超限）
    """

    def _split(t: str, sep_idx: int) -> List[str]:
        if len(t) <= max_size:
            return [t] if t else []
        if sep_idx >= len(SPLIT_SEPARATORS):
            return [t[i:i + max_size] for i in range(0, len(t), max_size)]

        sep = SPLIT_SEPARATORS[sep_idx]
        if sep == '':
            return [t[i:i + max_size] for i in range(0, len(t), max_size)]

        parts = t.split(sep)
        if len(parts) == 1:
            return _split(t, sep_idx + 1)

        result: List[str] = []
        for i, part in enumerate(parts):
            frag = part + sep if i < len(parts) - 1 else part
            if not frag.strip():
                continue
            if len(frag) <= max_size:
                result.append(frag)
            else:
                result.extend(_split(frag, sep_idx + 1))
        return result

    raw_chunks = _split(text, 0)
    if not raw_chunks:
        return []

    # 贪心合并：只要累积长度 + 下一片段 ≤ max_size，持续合并
    merged: List[str] = [raw_chunks[0]]
    for chunk in raw_chunks[1:]:
        prev = merged[-1]
        if len(prev) + len(chunk) <= max_size:
            merged[-1] = prev + chunk
        else:
            merged.append(chunk)

    # 兜底：末尾孤立的过小片段（< min_size）合并到前一个
    if len(merged) >= 2 and len(merged[-1]) < min_size:
        merged[-2] = merged[-2] + merged[-1]
        merged.pop()

    return merged


def merge_short_sections(
    sections: list[dict],
    min_size: int,
    max_size: int,
) -> list[dict]:
    """
    合并过短的 section，减少碎片化 chunk。

    算法：向前贪心 + 链式吸收
      - 正向遍历，短 section（meaningful chars < min_size）优先向前（next）合并，
        其次向后（prev）合并；合并后不推进索引，实现链式吸收。
      - 合并上限：两段 text 拼接后长度（含 \\n）不超过 max_size。
      - 两方向均超限时保留原样（安全阀）。

    元数据继承规则：
      - 向前合并（cur → next）：next 继承 cur 的 section_id 和 page_or_sheet
      - 向后合并（cur → prev）：prev 保持自己的 section_id 和 page_or_sheet
      - section_title：取两者中第一个非空值

    不修改原列表和原 dict。
    """
    if not sections:
        return []

    # 浅拷贝列表，每个 dict 也拷贝一份
    buf = [dict(s) for s in sections]

    i = 0
    while i < len(buf):
        cur = buf[i]
        if count_meaningful_chars(cur["text"]) >= min_size:
            i += 1
            continue

        # cur 是短 section，尝试合并
        merged = False

        # 尝试向前合并（与 next）
        if i + 1 < len(buf):
            nxt = buf[i + 1]
            combined_len = len(cur["text"]) + len(nxt["text"]) + 1  # +1 for \n
            if combined_len <= max_size:
                # 将 cur 文本 prepend 到 next
                nxt["text"] = cur["text"] + "\n" + nxt["text"]
                # next 继承 cur 的 section_id 和 page_or_sheet
                nxt["section_id"] = cur["section_id"]
                nxt["page_or_sheet"] = cur["page_or_sheet"]
                # section_title：取第一个非空值
                cur_title = cur.get("section_title", "")
                nxt_title = nxt.get("section_title", "")
                nxt["section_title"] = cur_title or nxt_title
                buf.pop(i)
                merged = True

        # 尝试向后合并（与 prev）
        if not merged and i > 0:
            prev = buf[i - 1]
            combined_len = len(prev["text"]) + len(cur["text"]) + 1  # +1 for \n
            if combined_len <= max_size:
                # 将 cur 文本 append 到 prev
                prev["text"] = prev["text"] + "\n" + cur["text"]
                # prev 保持自己的 section_id 和 page_or_sheet
                # section_title：取第一个非空值
                prev_title = prev.get("section_title", "")
                cur_title = cur.get("section_title", "")
                prev["section_title"] = prev_title or cur_title
                buf.pop(i)
                merged = True

        if not merged:
            # 安全阀：两方向都超 max_size，保留 cur
            i += 1

    return buf


# ==============================================================================
# 2-0c  表格处理优化（Phase 1：结构分离 + 格式转换）
# ==============================================================================

def _split_into_segments(raw_text: str, table_html_list: list) -> list:
    """
    将 Section 文本按表格位置分离为正文和表格片段序列。

    参数:
        raw_text: 原始 Section 文本（含 HTML 表格）
        table_html_list: HTML 表格列表（按出现顺序）

    返回:
        segments 列表，每个元素为：
        - {"type": "content", "text": "正文文本"}
        - {"type": "table", "text": "Markdown表格", "html": "<table>..."}
    """
    if not table_html_list:
        # 无表格，整个 section 作为正文
        text = raw_text.strip()
        if text:
            return [{"type": "content", "text": text}]
        return []

    segments = []
    last_end = 0

    for table_html in table_html_list:
        start = raw_text.find(table_html, last_end)
        if start == -1:
            continue

        # 表格前的正文
        if start > last_end:
            content_text = raw_text[last_end:start].strip()
            if content_text:
                segments.append({"type": "content", "text": content_text})

        # 表格本身（转换为 Markdown）
        table_markdown = _convert_single_table_to_markdown(table_html)
        segments.append({
            "type": "table",
            "text": table_markdown,   # Markdown 格式（用于 Embedding/Reranker/LLM）
            "html": table_html,       # 原始 HTML（仅 Web 渲染用）
        })

        last_end = start + len(table_html)

    # 最后一个表格后的正文
    if last_end < len(raw_text):
        content_text = raw_text[last_end:].strip()
        if content_text:
            segments.append({"type": "content", "text": content_text})

    return segments


def _convert_single_table_to_markdown(table_html: str) -> str:
    """
    将单个 HTML 表格转换为 Markdown。

    参数:
        table_html: 完整的 <table>...</table> HTML 字符串

    返回:
        Markdown 格式的表格字符串，转换失败时返回原 HTML
    """
    try:
        import pandas as pd
        from io import StringIO
        dfs = pd.read_html(StringIO(table_html), header=0)
        if dfs:
            df = dfs[0].fillna('').astype(str).replace('nan', '')
            return df.to_markdown(index=False)
    except Exception:
        pass
    return table_html  # 转换失败返回原文


def split_table_by_rows(table_html: str, max_rows: int = TABLE_MAX_ROWS) -> list:
    """
    按行数拆分表格，每个分块都保留表头。

    参数:
        table_html: 原始 HTML 表格
        max_rows: 每个分块最大数据行数（不含表头）

    返回:
        拆分后的表格信息列表，每个元素包含:
        - table_html: 拆分后的 HTML（含表头）
        - table_markdown: 拆分后的 Markdown（含表头）
        - row_range: (start, end) 行号范围
        - row_count: 数据行数
    """
    try:
        import pandas as pd
        from io import StringIO
        dfs = pd.read_html(StringIO(table_html), header=0)
        if not dfs:
            raise ValueError("No tables found")
        df = dfs[0]
    except Exception:
        # 解析失败，返回原表格
        return [{
            "table_html": table_html,
            "table_markdown": _convert_single_table_to_markdown(table_html),
            "row_range": (0, 0),
            "row_count": 0,
        }]

    total_rows = len(df)

    # 行数 ≤ max_rows，不拆分
    if total_rows <= max_rows:
        md = df.fillna('').astype(str).replace('nan', '').to_markdown(index=False)
        return [{
            "table_html": table_html,
            "table_markdown": md,
            "row_range": (0, total_rows),
            "row_count": total_rows,
        }]

    # 行数 > max_rows，按 max_rows 拆分
    results = []
    for start in range(0, total_rows, max_rows):
        end = min(start + max_rows, total_rows)
        sub_df = df.iloc[start:end]

        sub_md = sub_df.fillna('').astype(str).replace('nan', '').to_markdown(index=False)
        sub_html = sub_df.to_html(index=False, border=1)

        results.append({
            "table_html": sub_html,
            "table_markdown": sub_md,
            "row_range": (start, end),
            "row_count": end - start,
        })

    return results


def preprocess_section(section: dict) -> dict:
    """
    预处理 Section：转换格式 + 分离表格/正文。

    输入:
        section: {"section_id": "s0", "text": "...(含HTML表格)...", ...}

    输出:
        {
            "section_id": "s0",
            "page_or_sheet": "1",
            "section_title": "章节标题",
            "parent_text": "...(Markdown格式)...",  # 用于上下文
            "segments": [
                {"type": "content", "text": "正文A"},
                {"type": "table", "text": "Markdown表格", "html": "<table>..."},
                {"type": "content", "text": "正文B"},
                ...
            ]
        }
    """
    raw_text = section["text"]

    # 1. 提取所有 HTML 表格
    _, table_html_list = html_table_to_markdown(raw_text)

    # 2. 整体转换为 Markdown（用于 parent_text）
    parent_text, _ = html_table_to_markdown(raw_text)

    # 3. 分离表格和正文片段
    segments = _split_into_segments(raw_text, table_html_list)

    return {
        "section_id": section["section_id"],
        "page_or_sheet": section.get("page_or_sheet", "1"),
        "section_title": section.get("section_title", ""),
        "parent_text": parent_text,
        "segments": segments,
    }


def _merge_short_content_chunks(raw_chunks: list, min_size: int) -> list:
    """
    合并短正文片段。

    规则：
    1. 优先合并到上一个同类型（正文）chunk
    2. 不存在则合并到下一个同类型
    3. 前后都是表格则单独保留

    参数:
        raw_chunks: 原始 chunk 列表（含 type 字段）
        min_size: 短正文阈值（字符数）

    返回:
        合并后的 chunk 列表
    """
    if not raw_chunks:
        return raw_chunks

    result = []
    pending_merge = []  # 待向后合并的短正文

    for i, chunk in enumerate(raw_chunks):
        # 先处理待向后合并的内容
        if pending_merge and chunk["type"] == "content":
            for pm in pending_merge:
                chunk["text"] = pm["text"] + "\n" + chunk["text"]
                chunk["char_count"] = chunk.get("char_count", 0) + pm.get("char_count", 0) + 1
            pending_merge = []

        # 非正文或不短，直接加入
        if chunk["type"] != "content" or chunk.get("char_count", len(chunk["text"])) >= min_size:
            result.append(chunk)
            continue

        # 短正文，尝试向上合并
        merged = False
        for j in range(len(result) - 1, -1, -1):
            if result[j]["type"] == "content":
                result[j]["text"] += "\n" + chunk["text"]
                result[j]["char_count"] = result[j].get("char_count", 0) + chunk.get("char_count", 0) + 1
                merged = True
                break

        if not merged:
            # 无法向上合并，标记待向后合并
            pending_merge.append(chunk)

    # 处理末尾未合并的短正文（前后都是表格）
    for pm in pending_merge:
        result.append(pm)

    return result


def make_chunks_from_preprocessed_section(
    preprocessed: dict,
    file_record: dict,
    rel_path_normalized: str,
    max_size: int = 1200,
    min_size: int = 100,
    max_rows: int = TABLE_MAX_ROWS,
) -> list:
    """
    从预处理后的 Section 生成 chunks。

    参数:
        preprocessed: preprocess_section() 的输出
        file_record: 文件记录
        rel_path_normalized: 归一化后的相对路径
        max_size: 正文分块最大字符数
        min_size: 短正文合并阈值
        max_rows: 表格分块最大行数

    返回:
        chunks 列表，每个 chunk 包含完整字段
    """
    chunks = []
    content_index = 0
    table_index = 0

    section_id = preprocessed["section_id"]
    parent_id = f"{rel_path_normalized}#{section_id}"

    # 先收集所有片段的初步 chunks
    raw_chunks = []

    for segment in preprocessed["segments"]:
        if segment["type"] == "content":
            # 正文按字符数分块
            sub_texts = recursive_split(segment["text"], max_size, min_size)
            if not sub_texts:
                sub_texts = [segment["text"]] if segment["text"].strip() else []

            for i, sub_text in enumerate(sub_texts):
                if not sub_text.strip():
                    continue
                chunk_id = f"{rel_path_normalized}#{section_id}#c{content_index}"
                if len(sub_texts) > 1:
                    chunk_id += f"p{i}"
                raw_chunks.append({
                    "type": "content",
                    "chunk_id": chunk_id,
                    "text": sub_text,
                    "char_count": count_meaningful_chars(sub_text),
                })
            content_index += 1

        else:  # table
            # 表格按行数分块
            table_parts = split_table_by_rows(segment["html"], max_rows)
            for i, part in enumerate(table_parts):
                chunk_id = f"{rel_path_normalized}#{section_id}#t{table_index}"
                if len(table_parts) > 1:
                    chunk_id += f"p{i}"
                raw_chunks.append({
                    "type": "table",
                    "chunk_id": chunk_id,
                    "text": part["table_markdown"],
                    "table_markdown": part["table_markdown"],
                    "table_html": part["table_html"],
                    "table_rows": part["row_count"],
                    "char_count": count_meaningful_chars(part["table_markdown"]),
                })
            table_index += 1

    # 短正文合并
    raw_chunks = _merge_short_content_chunks(raw_chunks, min_size)

    # 组装最终 chunks
    for rc in raw_chunks:
        chunk = {
            "chunk_id": rc["chunk_id"],
            "parent_id": parent_id,
            "file_path": file_record["file_path"],
            "file_name": file_record["file_name"],
            "folder_code": file_record.get("folder_code"),
            "page_or_sheet": preprocessed["page_or_sheet"],
            "section_title": preprocessed["section_title"],
            "text": rc["text"],
            "char_count": rc.get("char_count", count_meaningful_chars(rc["text"])),
        }

        if rc["type"] == "table":
            chunk["is_table"] = True
            chunk["table_markdown"] = rc["table_markdown"]
            chunk["table_html"] = rc["table_html"]
            chunk["table_rows"] = rc["table_rows"]

        chunks.append(chunk)

    return chunks


def make_chunks_from_sections(
    sections: list,
    file_record: dict,
    max_size: int,
    min_size: int,
    max_rows: int = TABLE_MAX_ROWS,
) -> dict:
    """
    将 sections 列表展开为 ChunkRecord 字典。

    Phase 1 表格处理优化：
    - 表格与正文分离，分别用不同策略分块
    - 正文按字符数分块（max_size=1200）
    - 表格按行数分块（max_rows=30），每块补表头
    - parent_text 独立存储，避免冗余（节省 ~60MB）

    参数:
        sections: section 列表，格式：
            [{"section_id": str, "page_or_sheet": str, "text": str,
              "section_title": str (optional)}, ...]
        file_record: 文件记录
        max_size: 正文分块最大字符数
        min_size: 短正文合并阈值
        max_rows: 表格分块最大行数

    返回:
        {
            "parents": {parent_id: parent_text, ...},
            "chunks": [chunk_dict, ...]
        }
    """
    # 1. Section 级别合并（保持原有逻辑）
    sections = merge_short_sections(sections, min_size=min_size, max_size=max_size)

    rel_path = file_record.get("relative_path", file_record["file_name"])
    rel_path_normalized = rel_path.replace(os.sep, "/")

    parents = {}
    chunks = []

    for section in sections:
        # 2. Section 预处理（HTML→Markdown + 分离表格/正文）
        preprocessed = preprocess_section(section)
        parent_id = f"{rel_path_normalized}#{preprocessed['section_id']}"

        # 3. 存储 parent_text（独立存储，不重复）
        if preprocessed["parent_text"].strip():
            parents[parent_id] = preprocessed["parent_text"]

        # 4. 跳过空 section
        if not preprocessed["segments"]:
            continue

        # 5. 生成 chunks
        section_chunks = make_chunks_from_preprocessed_section(
            preprocessed, file_record, rel_path_normalized,
            max_size, min_size, max_rows
        )
        chunks.extend(section_chunks)

    return {"parents": parents, "chunks": chunks}


# ==============================================================================
# 2-1  GLM-OCR 调用封装
# ==============================================================================

# 智谱 OCR 并发控制（Semaphore 惰性初始化，线程安全）
_zhipu_ocr_semaphore: threading.Semaphore | None = None
_zhipu_sem_lock = threading.Lock()

def _get_zhipu_semaphore() -> threading.Semaphore:
    global _zhipu_ocr_semaphore
    if _zhipu_ocr_semaphore is None:
        with _zhipu_sem_lock:
            if _zhipu_ocr_semaphore is None:
                from config import ZHIPU_OCR_CONCURRENCY
                _zhipu_ocr_semaphore = threading.Semaphore(ZHIPU_OCR_CONCURRENCY)
    return _zhipu_ocr_semaphore


def _call_zhipu_layout_parsing(
    file_data: bytes,
    file_type: str = "image",   # "image" 或 "pdf"
    start_page: int | None = None,
    end_page: int | None = None,
) -> dict:
    """
    调用智谱线上 layout_parsing API。

    参数:
        file_data: 文件二进制数据（PNG/JPEG/PDF）
        file_type: "image" 或 "pdf"
        start_page/end_page: PDF 页码范围（可选）

    返回:
        {"md_results": str, "layout_details": list}
        或抛出 RuntimeError / requests.HTTPError

    注意:
        - 单图 ≤ 10MB，PDF ≤ 50MB，最多 100 页
        - layout_details 中 bbox_2d 为像素坐标，region 含 width/height 页面尺寸
    """
    import requests
    import base64
    from config import ZHIPU_API_KEY, ZHIPU_API_BASE_URL, GLM_OCR_MODEL

    if not ZHIPU_API_KEY:
        raise RuntimeError("ZHIPU_API_KEY 未配置，无法调用智谱 layout_parsing API")

    b64 = base64.b64encode(file_data).decode()

    # 智谱 API 要求 data URI 格式（非 base64:// 前缀）
    if file_type == "pdf":
        data_uri = f"data:application/pdf;base64,{b64}"
    else:
        # 自动检测 PNG / JPEG
        if file_data[:4] == b'\x89PNG':
            mime = "image/png"
        elif file_data[:2] == b'\xff\xd8':
            mime = "image/jpeg"
        else:
            mime = "image/png"  # 默认 PNG
        data_uri = f"data:{mime};base64,{b64}"

    payload = {
        "model": GLM_OCR_MODEL,
        "file": data_uri,
    }
    if file_type == "pdf":
        if start_page is not None:
            payload["start_page_id"] = start_page
        if end_page is not None:
            payload["end_page_id"] = end_page

    url = f"{ZHIPU_API_BASE_URL}/layout_parsing"
    headers = {
        "Authorization": f"Bearer {ZHIPU_API_KEY}",
        "Content-Type": "application/json",
    }

    # 并发控制：限制同时访问智谱 API 的线程数
    with _get_zhipu_semaphore():
        resp = requests.post(url, headers=headers, json=payload, timeout=300)
    resp.raise_for_status()

    result = resp.json()
    # 智谱返回结构：{"data": {"md_results": ..., "layout_details": ...}, ...}
    data = result.get("data", result)
    return {
        "md_results": data.get("md_results", ""),
        "layout_details": data.get("layout_details", []),
    }


# ==============================================================================
# 2-1b  VLM 图片分类与描述调用封装（DashScope qwen3-vl-plus）
# ==============================================================================

_vlm_context:  dict = {}  # {code: enhanced_query_text}，由 configure_vlm_context() 注入
_vlm_cache:    dict = {}  # {sha256_hex: {"type": str, "description": str}}
_vlm_cache_lock = threading.Lock()  # 文件级并发提取时保护 _vlm_cache 读写


def configure_vlm_context(enhanced_queries: dict) -> None:
    """
    注入 {code: enhanced_query_text} 映射，供 VLM prompt 提供上下文。
    在 align_evidence.py 加载增强查询文本后调用。
    """
    global _vlm_context
    _vlm_context = enhanced_queries or {}


def load_vlm_cache(cache_path: str | None = None) -> None:
    """从 cache_path（默认 VLM_CACHE_PATH）加载已有缓存到 _vlm_cache。align_evidence.py 启动时调用。"""
    global _vlm_cache
    import json as _json
    path = cache_path if cache_path is not None else VLM_CACHE_PATH
    if os.path.isfile(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                _vlm_cache = _json.load(f)
            print(f"  ✓ VLM 缓存加载：{len(_vlm_cache)} 条")
        except Exception:
            _vlm_cache = {}


def save_vlm_cache(cache_path: str | None = None) -> None:
    """将 _vlm_cache 持久化到 cache_path（默认 VLM_CACHE_PATH）。align_evidence.py 阶段 2a 完成后调用。"""
    import json as _json
    if not _vlm_cache:
        return
    path = cache_path if cache_path is not None else VLM_CACHE_PATH
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            _json.dump(_vlm_cache, f, ensure_ascii=False, indent=2)
        print(f"  ✓ VLM 缓存已保存：{len(_vlm_cache)} 条 → {path}")
    except Exception as e:
        print(f"  [警告] VLM 缓存保存失败：{e}")


# 合法图片类型集合（VLM 分类输出选项）
_VLM_VALID_TYPES = {"文档扫描件", "表格截图", "照片", "流程图", "数据图表", "证书", "其他"}


def _parse_vlm_response(raw: str) -> dict:
    """
    解析 VLM 返回的分类+描述文本。

    期望格式：
        类型：xxx
        描述：yyy

    容错：若解析失败，type 设为 "其他"，description 取原始输出前 100 字。
    """
    result = {"type": "其他", "description": ""}

    type_match = re.search(r'类型[：:]\s*(.+?)(?:\n|$)', raw)
    desc_match = re.search(r'描述[：:]\s*(.+?)(?:\n|$)', raw)

    if type_match:
        t = type_match.group(1).strip()
        if t in _VLM_VALID_TYPES:
            result["type"] = t
        else:
            # 模糊匹配：包含关系
            for valid in _VLM_VALID_TYPES:
                if valid in t or t in valid:
                    result["type"] = valid
                    break

    if desc_match:
        result["description"] = desc_match.group(1).strip()[:200]
    elif raw.strip():
        # 回退：使用原始输出截断
        result["description"] = raw.strip()[:100]

    return result


def call_vlm_classify(
    img_bytes: bytes,
    filename: str = "",
    page: str = "",
    idx: int = 0,
    folder_code: str | None = None,
    source_path: str = "",
) -> dict | None:
    """
    调用 DashScope qwen3-vl-plus 对图片进行分类和描述。

    参数:
        img_bytes    - PNG 格式图片字节流（应已 resize）
        filename     - 来源文件名（供 prompt 上下文）
        page         - 所在页码/sheet
        idx          - 该页/文件中的第几张图（0-based）
        folder_code  - 路径编码（若有）
        source_path  - 来源文件相对路径（供缓存溯源，如 G-公司治理/GA1/报告.pdf）

    返回:
        {"type": str, "description": str, "source": str, "page": str, "idx": int}
        或 None（调用失败时）

    分类选项：文档扫描件 / 表格截图 / 照片 / 流程图 / 数据图表 / 证书 / 其他
    """
    import base64
    import hashlib
    from dashscope import MultiModalConversation
    from config import DASHSCOPE_API_KEY

    # ── 缓存命中检查 ─────────────────────────────────────────────────────
    img_hash = hashlib.sha256(img_bytes).hexdigest()
    with _vlm_cache_lock:
        if img_hash in _vlm_cache:
            return _vlm_cache[img_hash]

    # ── 构建 prompt ───────────────────────────────────────────────────────
    location_str = f"第{page}页 第{idx+1}张图" if page else f"第{idx+1}张图"
    context_parts = []
    if filename:
        context_parts.append(f"- 文件：{filename}")
    context_parts.append(f"- 位置：{location_str}")

    if folder_code and folder_code in _vlm_context:
        context_parts.append(f"- 所属分类：{folder_code} — {_vlm_context[folder_code]}")
    elif folder_code:
        context_parts.append(f"- 所属分类：{folder_code}")

    context_block = "\n".join(context_parts)

    prompt = (
        "你是一个图片分类与描述助手。请根据以下信息对图片进行分类和简要描述。\n\n"
        f"图片来源信息：\n{context_block}\n\n"
        "请输出：\n"
        "1. 图片类型（从以下选一个）：文档扫描件 / 表格截图 / 照片 / 流程图 / "
        "数据图表 / 证书 / 其他\n"
        "2. 一句话描述（30-80字，说清楚图片展示了什么内容）\n\n"
        "输出格式（严格遵守，不要多余文字）：\n"
        "类型：xxx\n"
        "描述：yyy"
    )

    # ── 调用 DashScope VLM（带重试） ────────────────────────────────────────
    import time as _time
    from config import API_MAX_RETRIES, API_RETRY_BASE_DELAY

    data_uri = "data:image/png;base64," + base64.b64encode(img_bytes).decode()

    for attempt in range(API_MAX_RETRIES):
        try:
            response = MultiModalConversation.call(
                api_key=DASHSCOPE_API_KEY,
                model=VLM_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"image": data_uri},
                        {"text": prompt},
                    ],
                }],
                enable_thinking=False,
            )

            # 提取文本
            raw = ""
            if hasattr(response, 'output') and response.output:
                choices = response.output.get('choices', [])
                if choices:
                    content = choices[0].get('message', {}).get('content', [])
                    if isinstance(content, list) and content:
                        raw = content[0].get('text', '')
                    elif isinstance(content, str):
                        raw = content

            if not raw:
                if attempt < API_MAX_RETRIES - 1:
                    wait = API_RETRY_BASE_DELAY * (2 ** attempt)
                    print(f"  [警告] VLM 返回空内容（第 {attempt + 1} 次）：{filename} p{page}，"
                          f"等待 {wait:.0f}s 后重试")
                    _time.sleep(wait)
                    continue
                print(f"  [警告] VLM 返回空内容（已重试 {API_MAX_RETRIES} 次）：{filename} p{page}")
                return None

            result = _parse_vlm_response(raw)
            result["source"] = source_path or filename
            result["page"]   = page
            result["idx"]    = idx
            with _vlm_cache_lock:
                _vlm_cache[img_hash] = result
            return result

        except Exception as e:
            wait = API_RETRY_BASE_DELAY * (2 ** attempt)
            if attempt < API_MAX_RETRIES - 1:
                print(f"  [警告] VLM 调用第 {attempt + 1} 次失败：{e}，"
                      f"等待 {wait:.0f}s 后重试")
                _time.sleep(wait)
            else:
                print(f"  [警告] VLM 调用失败（已重试 {API_MAX_RETRIES} 次）：{e}")
                return None


# ==============================================================================
# 2-1c  图片过滤 / Resize / 格式转换 工具函数
# ==============================================================================

def _filter_image(img_bytes: bytes, width: int, height: int) -> bool:
    """
    判断图片是否通过最小尺寸和最小字节过滤。
    通过 → True（应处理），不通过 → False（应跳过）。
    """
    if len(img_bytes) < IMAGE_MIN_BYTES:
        return False
    if width < IMAGE_MIN_WIDTH or height < IMAGE_MIN_HEIGHT:
        return False
    return True


def _resize_image_bytes(img_bytes: bytes, max_px: int | None = None) -> bytes:
    """
    将图片字节流等比缩放至长边不超过 max_px，返回 PNG 字节流。
    若已在范围内或 max_px 为 None，直接返回原始字节。
    """
    from PIL import Image
    import io

    if max_px is None:
        max_px = VLM_MAX_IMAGE_PX

    img = Image.open(io.BytesIO(img_bytes))
    w, h = img.size
    if max(w, h) <= max_px:
        return img_bytes

    scale = max_px / max(w, h)
    new_w = int(w * scale)
    new_h = int(h * scale)
    img = img.resize((new_w, new_h), Image.LANCZOS)

    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="PNG")
    return buf.getvalue()


def _to_png_bytes(img_bytes: bytes, ext: str = "") -> tuple[bytes, int, int]:
    """
    将任意图片格式（JPEG/PNG/BMP/TIFF 等）转为 PNG 字节流。
    返回 (png_bytes, width, height)。
    """
    from PIL import Image
    import io

    img = Image.open(io.BytesIO(img_bytes))
    w, h = img.size
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="PNG")
    return buf.getvalue(), w, h


def _assemble_image_section_text(
    vlm_result: dict,
    filename: str,
    page: str,
    idx: int,
    png_bytes: bytes | None = None,
) -> str:
    """
    根据 VLM 分类结果组装 section text。

    对于 文档扫描件/表格截图：调用 SDK 进行版面检测 + 结构化 OCR。
    对于其他类型：使用 VLM 描述。

    格式：
        [图片] 来源：{filename} 第{page}页 | 类型：{type}
        {description_or_ocr_text}

    注意（冗余兜底逻辑）:
        在 SDK 路径（parse_pdf_sdk）中，表格区域已被 PP-DocLayout-V3 识别并
        用 "Table Recognition:" prompt OCR，VLM 看到的 image/chart 区域
        通常不含表格。此处的表格 OCR 逻辑在 SDK 路径下理论上冗余，但作为
        边界兜底（Paddle 可能漏判低质量表格截图）仍保留，且对 PyMuPDF 路径
        和 DOCX/独立图片等场景仍有实际作用。
    """
    img_type    = vlm_result.get("type", "其他")
    description = vlm_result.get("description", "")

    # 文档扫描件/表格截图 → 使用 SDK 进行版面检测 + 结构化 OCR
    final_text = description
    if img_type in ("文档扫描件", "表格截图") and png_bytes is not None:
        try:
            ocr_text = ocr_image_with_sdk(png_bytes, img_type=img_type)
            if ocr_text and count_meaningful_chars(ocr_text) > 10:
                final_text = ocr_text  # OCR 文本替代 VLM 描述
        except Exception:
            pass  # OCR 失败时保留 VLM 描述

    header = f"[图片] 来源：{filename} 第{page}页 | 类型：{img_type}"
    return f"{header}\n{final_text}"


def _process_single_image(
    png_bytes: bytes,
    width: int,
    height: int,
    filename: str,
    page: str,
    idx: int,
    folder_code: str | None = None,
    source_path: str = "",
) -> str | None:
    """
    图片处理公共管线：过滤 → 缩放 → VLM 分类+描述 → 组装 section 文本。

    参数:
        png_bytes    - 已转为 PNG 格式的图片字节流（原始尺寸，OCR 用）
        width/height - 图片原始尺寸（用于过滤判断）
        filename     - 来源文件名
        page         - 所在页码/sheet（字符串）
        idx          - 该页/文件中的图片序号（0-based，用于 VLM prompt）
        folder_code  - 路径编码（若有）
        source_path  - 来源文件相对路径（供缓存溯源）

    返回:
        组装好的 section 文本字符串，或 None（过滤/VLM 不可用时）。
    """
    if not _filter_image(png_bytes, width, height):
        return None

    png_resized = _resize_image_bytes(png_bytes)

    vlm_result = call_vlm_classify(
        img_bytes=png_resized,
        filename=filename,
        page=page,
        idx=idx,
        folder_code=folder_code,
        source_path=source_path,
    )

    if vlm_result is None:
        return None

    return _assemble_image_section_text(
        vlm_result=vlm_result,
        filename=filename,
        page=page,
        idx=idx + 1,
        png_bytes=png_bytes,
    )


def _process_images_batch(
    image_tasks: list,
    max_concurrent: int | None = None,
) -> list:
    """
    并发批量处理图片（VLM 分类+描述 → 可能追加 OCR）。

    参数:
        image_tasks: 字典列表，每个元素包含 _process_single_image() 的关键字参数：
            {"png_bytes", "width", "height", "filename", "page", "idx",
             "folder_code"(可选), "source_path"(可选)}
        max_concurrent: 最大并发数，默认从 config.VLM_CONCURRENCY 读取

    返回:
        与输入等长的列表，成功为 section 文本字符串，失败/过滤为 None。
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from config import VLM_CONCURRENCY

    if max_concurrent is None:
        max_concurrent = VLM_CONCURRENCY
    if not image_tasks:
        return []

    results = [None] * len(image_tasks)

    with ThreadPoolExecutor(max_workers=max_concurrent) as pool:
        future_to_idx = {
            pool.submit(_process_single_image, **task): i
            for i, task in enumerate(image_tasks)
        }
        for future in as_completed(future_to_idx):
            i = future_to_idx[future]
            try:
                results[i] = future.result()
            except Exception as e:
                print(f"  [警告] 图片并发处理异常: {e}")

    return results


def call_glmocr(img_bytes: bytes) -> str:
    """
    调用智谱线上 layout_parsing API 对单张图片进行 OCR。

    img_bytes: PNG/JPEG 格式的图片字节流
    返回: 识别出的文本（Markdown 格式）
    """
    result = _call_zhipu_layout_parsing(img_bytes, file_type="image")
    return result.get("md_results", "").strip()


def ocr_image_with_sdk(img_bytes: bytes, img_type: str = "文档扫描件") -> str:
    """
    对单张图片进行版面检测 + 结构化 OCR（智谱线上 layout_parsing API）。

    带重试：指数退避重试 API_MAX_RETRIES 次，全部失败后回退到 call_glmocr()。

    参数:
        img_bytes: PNG 格式图片字节流
        img_type: VLM 分类结果（用于日志，不影响行为）

    返回:
        识别出的文本（Markdown 格式，表格为 HTML）
    """
    import time as _time
    from config import API_MAX_RETRIES, API_RETRY_BASE_DELAY

    for attempt in range(API_MAX_RETRIES):
        try:
            result = _call_zhipu_layout_parsing(img_bytes, file_type="image")
            markdown = result.get("md_results", "").strip()
            if markdown:
                return markdown
            # md_results 为空，视为失败触发重试
            if attempt < API_MAX_RETRIES - 1:
                wait = API_RETRY_BASE_DELAY * (2 ** attempt)
                _time.sleep(wait)
                continue
        except Exception as e:
            wait = API_RETRY_BASE_DELAY * (2 ** attempt)
            if attempt < API_MAX_RETRIES - 1:
                print(f"  [警告] 图片 OCR 第 {attempt + 1} 次失败：{e}，"
                      f"等待 {wait:.0f}s 后重试")
                _time.sleep(wait)
            else:
                print(f"  [警告] 图片 OCR 失败（已重试 {API_MAX_RETRIES} 次），回退到 call_glmocr：{e}")

    return call_glmocr(img_bytes)


# ==============================================================================
# 2-1  PDF 提取
# ==============================================================================


# ==============================================================================
# 2-1v2  PDF v2 路径：智谱线上 layout_parsing API
# ==============================================================================


def classify_pdf_v2(doc) -> str:
    """
    v2 整文档级 PDF 分流，返回 "pymupdf" 或 "sdk"。

    分流逻辑：
      1. PPT 判断（复用现有 PPT 检测逻辑）→ 命中返回 "sdk"
      2. 逐页统计有效字符数
      3. 存在任何一页 < PDF_PAGE_MIN_CHARS → 返回 "sdk"（混合 PDF / 扫描件）
      4. 全部页 >= PDF_PAGE_MIN_CHARS → 返回 "pymupdf"（纯文字 PDF）

    与 v1 classify_pdf() 的区别：
      - v1 用页均字符判断，会漏掉混合 PDF
      - v2 逐页检测，任一页低于阈值即走 SDK
    """
    total_pages = len(doc)
    if total_pages == 0:
        return "pymupdf"

    # ── PPT 转 PDF 判断（复用现有逻辑）──────────────────────────────────────
    sample_pages = [doc[i] for i in range(min(PDF_SAMPLE_PAGES, total_pages))]

    aspect_ratios   = []
    avg_block_chars = []
    img_counts      = []

    for page in sample_pages:
        rect = page.rect
        aspect_ratios.append(rect.width / rect.height if rect.height > 0 else 1.0)

        blocks      = page.get_text("blocks")
        text_blocks = [b for b in blocks if b[6] == 0]  # type 0 = text
        if text_blocks:
            avg_block_chars.append(sum(len(b[4]) for b in text_blocks) / len(text_blocks))
        else:
            avg_block_chars.append(0)

        img_counts.append(len(page.get_images(full=False)))

    mean_aspect    = sum(aspect_ratios)   / len(aspect_ratios)
    mean_blk_chars = sum(avg_block_chars) / len(avg_block_chars)
    mean_imgs      = sum(img_counts)      / len(img_counts)

    if mean_aspect > PDF_PPT_ASPECT_RATIO and (mean_blk_chars < PDF_PPT_AVG_BLOCK_CHARS or mean_imgs >= PDF_PPT_AVG_IMAGES):
        return "sdk"

    # ── 逐页检测有效字符（v2 核心改进）──────────────────────────────────────
    for page in doc:
        page_chars = count_meaningful_chars(page.get_text())
        if page_chars < PDF_PAGE_MIN_CHARS:
            return "sdk"

    return "pymupdf"


def _rebuild_title_levels_rule(titles: list) -> list:
    """
    规则方案：用编号模式推断标题层级。

    输入:
        [{"index": 0, "sdk_label": "doc_title"|"paragraph_title",
          "raw_text": "1.1 目的"}, ...]

    输出:
        [{"index": 0, "level": 2, "text": "1.1 目的"}, ...]

    逻辑:
        - sdk_label == "doc_title" → level = 1
        - sdk_label == "paragraph_title" → 调用 _heading_numeric_level()
          - 返回 > 0 → 使用返回值
          - 返回 0（无编号）→ 默认 level = 2
    """
    result = []
    for t in titles:
        idx = t["index"]
        text = t["raw_text"]
        label = t["sdk_label"]

        if label == "doc_title":
            level = 1
        else:
            # paragraph_title → 用编号模式推断
            lv = _heading_numeric_level(text)
            level = lv if lv > 0 else 2

        result.append({"index": idx, "level": level, "text": text})

    return result


def _rebuild_title_levels_llm(titles: list, file_name: str = "") -> list:
    """
    LLM 方案：将标题序列发给 LLM 推断层级。

    输入/输出格式同 _rebuild_title_levels_rule()。
    file_name: 文件名，作为 LLM 上下文信息辅助判断。

    逻辑:
        1. 加载 src/prompts/title_level_rebuild.txt 模板
        2. 构建 JSON 标题列表 + 文件名上下文传入 LLM
        3. 解析返回的 JSON 数组 [{index, level}]
        4. 失败时回退到规则方案
    """
    import json as _json
    from pathlib import Path

    # 加载 prompt 模板
    prompt_path = Path(__file__).parent / "prompts" / "title_level_rebuild.txt"
    try:
        prompt_template = prompt_path.read_text(encoding="utf-8")
    except Exception as e:
        print(f"  [警告] 加载标题层级 prompt 失败：{e}，回退到规则方案")
        return _rebuild_title_levels_rule(titles)

    # 构建输入 JSON
    titles_for_llm = [
        {"index": t["index"], "sdk_label": t["sdk_label"], "raw_text": t["raw_text"]}
        for t in titles
    ]
    titles_json = _json.dumps(titles_for_llm, ensure_ascii=False, indent=2)
    prompt = prompt_template.replace("{titles_json}", titles_json)
    prompt = prompt.replace("{file_name}", file_name or "未知文件")

    # 调用 LLM
    try:
        from openai import OpenAI
        client = OpenAI(
            api_key=DRAFT_LLM_API_KEY,
            base_url=DRAFT_LLM_BASE_URL if DRAFT_LLM_BASE_URL.endswith("/v1")
                     else DRAFT_LLM_BASE_URL + "/v1",
        )
        extra_body = {"enable_thinking": True} if DRAFT_ENABLE_THINKING else {}
        response = client.chat.completions.create(
            model=DRAFT_LLM_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            max_tokens=4096,
            extra_body=extra_body,
        )
        raw_response = response.choices[0].message.content or ""

        # 解析 JSON（处理可能的 markdown 代码块包裹）
        json_str = raw_response.strip()
        if json_str.startswith("```"):
            # 去除 ```json ... ``` 包裹
            lines = json_str.split("\n")
            lines = [l for l in lines if not l.strip().startswith("```")]
            json_str = "\n".join(lines)

        levels = _json.loads(json_str)

        # 校验完整性
        if len(levels) != len(titles):
            print(f"  [警告] LLM 返回标题数 {len(levels)} ≠ 输入 {len(titles)}，回退到规则方案")
            return _rebuild_title_levels_rule(titles)

        # 构建结果
        result = []
        for item, orig in zip(levels, titles):
            level = int(item.get("level", 2))
            level = max(1, min(6, level))  # 限制范围 1-6
            result.append({
                "index": orig["index"],
                "level": level,
                "text": orig["raw_text"],
            })

        return result

    except Exception as e:
        print(f"  [警告] LLM 标题层级推断失败：{e}，回退到规则方案")
        return _rebuild_title_levels_rule(titles)


def _parse_sdk_markdown(markdown: str, titles_with_levels: list) -> list:
    """
    基于标题层级重建结果，将 SDK Markdown 切割为 section 列表。

    输入:
        markdown: SDK 输出的 Markdown 文本
        titles_with_levels: 标题层级重建结果
            [{"index": 0, "level": 2, "text": "1.1 目的"}, ...]

    输出:
        [{"section_id": "s0", "page_or_sheet": "1",
          "text": "...", "section_title": ""}, ...]

    算法:
        Step 1: 将 Markdown 中的 # / ## 标题行匹配到 titles_with_levels，
                替换为真实层级
        Step 2: 确定切割层级（复用 DOCX 的动态策略）
        Step 3: 按 cut_level 分割 section
    """
    lines = markdown.split("\n")

    # ── Step 1: 标题行匹配 + 层级替换 ────────────────────────────────────────
    # 构建标题文本到层级的映射（按顺序消费）
    title_queue = list(titles_with_levels)  # 复制，避免修改原列表
    title_idx = 0  # 下一个待匹配的标题索引

    # 标记每行：是否标题、层级
    line_records = []
    for line in lines:
        stripped = line.strip()
        # 匹配 # 或 ## 开头的标题行
        heading_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
        if heading_match and title_idx < len(title_queue):
            heading_text = heading_match.group(2).strip()
            expected_text = title_queue[title_idx]["text"].strip()

            # 模糊匹配：标题文本是否一致（SDK markdown 中标题文本应与 json_result 一致）
            if heading_text == expected_text or heading_text.startswith(expected_text[:10]):
                level = title_queue[title_idx]["level"]
                title_idx += 1
                line_records.append({
                    "is_title": True,
                    "level":    level,
                    "text":     heading_text,
                    "raw_line": line,
                })
                continue

        # 非标题行（普通文本、表格、空行等）
        line_records.append({
            "is_title": False,
            "level":    0,
            "text":     line,
            "raw_line": line,
        })

    # ── Step 2: 确定切割层级（复用 DOCX 动态策略） ───────────────────────────
    max_level = 0
    for rec in line_records:
        if rec["is_title"] and rec["level"] > max_level:
            max_level = rec["level"]

    # 无标题 → 整文件作一个 section
    if max_level == 0:
        full_text = "\n".join(r["text"] for r in line_records).strip()
        if not full_text:
            return []
        return [{
            "section_id":    "doc",
            "page_or_sheet": "1",
            "text":          full_text,
            "section_title": "",
        }]

    # max_level <= 2 → cut_level = 1
    # max_level >= 3 → cut_level = 2
    cut_level = 1 if max_level <= 2 else 2

    # ── Step 3: 按 cut_level 分割 section ────────────────────────────────────
    sections = []
    current_l1     = ""      # 第1层标题（仅 cut_level=2 时用作 section_title）
    current_title  = ""      # 当前 section 的标题行
    body_parts     = []
    counter        = 0

    def _flush():
        nonlocal counter
        body     = "\n".join(body_parts).strip()
        sec_text = (current_title + "\n" + body).strip() if current_title else body
        if sec_text:
            sections.append({
                "section_id":    f"s{counter}",
                "page_or_sheet": "1",  # SDK 不提供逐行页码，统一为 "1"
                "text":          sec_text,
                "section_title": current_l1 if cut_level == 2 else "",
            })
            counter += 1

    for rec in line_records:
        if not rec["is_title"]:
            body_parts.append(rec["text"])
            continue

        lv = rec["level"]

        if cut_level == 1:
            # 所有标题都触发切割
            _flush()
            current_title = rec["text"]
            body_parts    = []
        else:
            # cut_level == 2
            if lv == 1:
                # 第1层：不切，更新 l1 上下文；将标题文本并入当前 body
                _flush()
                current_l1    = rec["text"]
                current_title = ""
                body_parts    = [rec["text"]]  # l1 标题保留在正文开头
            elif lv == 2:
                # 第2层：触发切割
                _flush()
                current_title = rec["text"]
                body_parts    = []
            else:
                # 第3层及以下：不切，作为正文保留
                body_parts.append(rec["text"])

    _flush()

    return sections


def _extract_sdk_images(json_result, doc, file_record: dict, img_base_dir: str | None = None) -> list:
    """
    处理 SDK JSON 中 label=="image" 或 "chart" 的区域。

    输入:
        json_result: SDK 的 json_result（分页数组，每页含区域列表）
                     格式: [[{index, label, native_label, bbox_2d, content, ...}, ...], ...]
        doc: fitz.Document（用于按 bbox 裁切图片）
        file_record: 文件记录

    输出:
        image section 列表

    逻辑:
        1. 遍历 json_result 的每页每个区域
        2. 筛选 native_label 为 "image" 或 "chart" 的区域
        3. 用 bbox 坐标从对应页面裁切图片（page.get_pixmap(clip=rect)）
        4. 保存裁切图片到 img_base_dir（供后续 VLM 描述溯源）
        5. 调用 _process_single_image() → VLM 分类+描述
        6. 返回 image sections

    注意（冗余兜底逻辑）:
        在 SDK 路径中，PP-DocLayout-V3 已将表格区域识别为 `table` 类型并用
        "Table Recognition:" prompt 进行结构化 OCR，因此这里 VLM 看到的
        image/chart 区域**通常不包含表格**。后续 _process_single_image 中
        对"文档扫描件/表格截图"类型追加 OCR 的逻辑，在 SDK 路径下理论上冗余，
        但作为边界兜底（Paddle 可能漏判低质量表格截图）仍保留。
    """
    import fitz

    sections   = []
    filename    = file_record.get("file_name", "")
    folder_code = file_record.get("folder_code")
    source_path = file_record.get("relative_path", filename)
    global_img_idx = 0

    # 文件级图片缓存目录：{img_base_dir}/{文件名去扩展名}/
    file_stem = os.path.splitext(filename)[0]
    img_dir = os.path.join(img_base_dir, file_stem)

    if not json_result or not isinstance(json_result, list):
        return sections

    for page_idx, page_regions in enumerate(json_result):
        if page_idx >= len(doc):
            break  # json_result 页数超出 PDF 页数

        page = doc[page_idx]

        if not isinstance(page_regions, list):
            continue

        for region in page_regions:
            native_label = region.get("native_label", "")
            if native_label not in ("image", "chart"):
                continue

            # 获取 bbox 坐标
            bbox = region.get("bbox_2d")
            if not bbox or len(bbox) != 4:
                continue

            try:
                # bbox_2d 格式: [x0, y0, x1, y1]（像素坐标）
                # 需要将 bbox 坐标转换为 PyMuPDF 坐标系（72 DPI）
                # 线上 API 区域含 width/height 字段 → 用页面尺寸计算 scale
                region_page_width = region.get("width")
                if not region_page_width or region_page_width <= 0:
                    print(f"  [警告] SDK 图片区域缺少 width 字段（页{page_idx + 1}），跳过")
                    continue
                scale = page.rect.width / region_page_width
                rect = fitz.Rect(
                    bbox[0] * scale,
                    bbox[1] * scale,
                    bbox[2] * scale,
                    bbox[3] * scale,
                )

                # 裁切图片（使用较高 DPI 获取清晰图片）
                clip_dpi = 150  # 裁切用 DPI
                mat = fitz.Matrix(clip_dpi / 72.0, clip_dpi / 72.0)
                pix = page.get_pixmap(matrix=mat, clip=rect)
                png_bytes = pix.tobytes("png")
                w, h = pix.width, pix.height

                # 保存裁切图片到本地（供后续 VLM 描述溯源）
                img_filename = f"p{page_idx + 1}_{native_label}_{global_img_idx}.png"
                try:
                    os.makedirs(img_dir, exist_ok=True)
                    img_path = os.path.join(img_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(png_bytes)
                except Exception:
                    pass  # 保存失败不影响提取

                text = _process_single_image(
                    png_bytes, w, h,
                    filename=filename,
                    page=str(page_idx + 1),
                    idx=global_img_idx,
                    folder_code=folder_code,
                    source_path=source_path,
                )
                global_img_idx += 1

                if text is None:
                    continue

                sections.append({
                    "section_id":    f"img_sdk_p{page_idx + 1}_{global_img_idx}",
                    "page_or_sheet": str(page_idx + 1),
                    "text":          text,
                    "section_title": "",
                })

            except Exception as e:
                print(f"  [警告] SDK 图片区域处理失败（页{page_idx + 1}）：{e}")
                continue

    return sections


def parse_pdf_sdk(file_record: dict, img_base_dir: str | None = None) -> list:
    """
    PDF 提取入口（扫描件/混合 PDF），使用智谱线上 layout_parsing API。

    流程:
        1. 调用智谱线上 API → 获取 markdown + layout_details
        2. 从 markdown 提取标题列表（# → doc_title, ## → paragraph_title）
        3. 标题层级重建（rule 或 llm 模式）
        4. Markdown 切割为 section 列表
        5. 处理 layout_details 中的 image/chart 区域
        6. 合并文本 sections + image sections

    参数:
        file_record: 文件记录 dict
        img_base_dir: SDK 图片缓存目录（企业专属路径，None 时跳过图片提取）

    返回:
        section 列表
    """
    import fitz

    file_path = file_record["file_path"]
    file_name = file_record.get("file_name", "")

    # ── 1. 智谱线上 API 解析（带重试） ─────────────────────────────────────

    import time as _time
    from config import API_MAX_RETRIES, API_RETRY_BASE_DELAY

    with open(file_path, "rb") as f:
        pdf_bytes = f.read()

    api_result = None
    for attempt in range(API_MAX_RETRIES):
        try:
            t0 = _time.time()
            print(f"    智谱线上解析中：{file_name} ...", flush=True)
            api_result = _call_zhipu_layout_parsing(pdf_bytes, file_type="pdf")
            elapsed = _time.time() - t0
            print(f"    ✓ 线上解析完成（{elapsed:.1f}s）：{file_name}")
            break
        except Exception as e:
            wait = API_RETRY_BASE_DELAY * (2 ** attempt)
            if attempt < API_MAX_RETRIES - 1:
                print(f"  [警告] PDF 解析第 {attempt + 1} 次失败：{file_name} — {e}，"
                      f"等待 {wait:.0f}s 后重试")
                _time.sleep(wait)
            else:
                print(f"  [错误] PDF 解析失败（已重试 {API_MAX_RETRIES} 次）：{file_name} — {e}")
                return []

    markdown    = api_result.get("md_results", "")
    json_result = api_result.get("layout_details", None)

    if not markdown.strip():
        print(f"  [警告] 返回空 markdown：{file_name}")
        return []

    # ── 3. 提取标题列表 ──────────────────────────────────────────────────────
    # 扫描 markdown 中 # / ## 开头的行
    title_list = []
    title_index = 0
    for line in markdown.split("\n"):
        stripped = line.strip()
        heading_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
        if heading_match:
            level_markers = heading_match.group(1)
            heading_text = heading_match.group(2).strip()
            # # → doc_title, ## → paragraph_title
            if len(level_markers) == 1:
                sdk_label = "doc_title"
            else:
                sdk_label = "paragraph_title"

            title_list.append({
                "index":     title_index,
                "sdk_label": sdk_label,
                "raw_text":  heading_text,
            })
            title_index += 1

    # ── 4. 标题层级重建 ──────────────────────────────────────────────────────
    if title_list:
        if TITLE_REBUILD_MODE == "llm":
            titles_with_levels = _rebuild_title_levels_llm(title_list, file_name=file_name)
        else:
            titles_with_levels = _rebuild_title_levels_rule(title_list)
    else:
        titles_with_levels = []

    # ── 5. Markdown 切割为 section ───────────────────────────────────────────
    text_sections = _parse_sdk_markdown(markdown, titles_with_levels)

    # ── 6. 图片区域处理 ──────────────────────────────────────────────────────
    image_sections = []
    if json_result and img_base_dir is not None:
        try:
            doc = fitz.open(file_path)
            image_sections = _extract_sdk_images(json_result, doc, file_record, img_base_dir=img_base_dir)
        except Exception as e:
            print(f"  [警告] SDK 图片区域处理失败：{file_name} — {e}")

    return text_sections + image_sections


def _find_title_threshold(
    all_sizes: List[float],
    para_records: List[dict],
) -> float:
    """
    根据字号分布和 block 级命中率，选取最佳标题阈值。

    算法（逐级上探）：
    1. 从 P75 百分位开始，计算该阈值下标题 block 占比
    2. 若占比 > PDF_TITLE_RATIO_MAX（50%），说明阈值过低
    3. 收集所有去重字号并升序排列，从当前阈值向上逐级尝试
    4. 取第一个满足条件的阈值：
       - 命中 block 数 > PDF_TITLE_MIN_COUNT（2）
       - 命中 block 占比 ≤ PDF_TITLE_RATIO_MAX（50%）
    5. 若所有更高字号均不满足 → 返回 0（放弃标题切割）

    Args:
        all_sizes:    全文所有 span 的 font size 列表（含重复）
        para_records: block 级记录列表，每条需含 'max_size' 和 'text' 字段

    Returns:
        float: 标题阈值（size ≥ 此值的 block 视为标题）。
               返回 0 表示无法识别有效标题层级。
    """
    if not all_sizes or not para_records:
        return 0

    # Step 1: P75 基础阈值（与原逻辑一致）
    sorted_sizes = sorted(all_sizes)
    p_idx = int(len(sorted_sizes) * PDF_TITLE_SIZE_PERCENTILE)
    base_threshold = sorted_sizes[min(p_idx, len(sorted_sizes) - 1)]

    # Step 2: 计算 block 级命中率
    def _title_stats(threshold: float) -> tuple:
        """返回 (命中数, 总数, 占比)"""
        total = len(para_records)
        hits = sum(1 for r in para_records
                   if r["max_size"] >= threshold
                   and len(r["text"]) <= PDF_TITLE_MAX_CHARS)
        return hits, total, hits / total if total > 0 else 0.0

    hits, total, ratio = _title_stats(base_threshold)

    # P75 命中率合理 → 直接采用
    if ratio <= PDF_TITLE_RATIO_MAX:
        return base_threshold if hits > 0 else 0

    # Step 3: P75 命中率过高 → 逐级上探
    unique_sizes = sorted(set(all_sizes))
    higher_levels = [s for s in unique_sizes if s > base_threshold]

    for level in higher_levels:
        hits, total, ratio = _title_stats(level)
        if hits > PDF_TITLE_MIN_COUNT and ratio <= PDF_TITLE_RATIO_MAX:
            return level

    # Step 4: 所有更高字号都不满足 → 放弃标题切割
    return 0


def parse_normal_pdf(doc) -> List[dict]:
    """
    对普通文字 PDF 按标题结构切分 section，返回 section 列表。

    标题识别逻辑（逐级上探）：
    1. 收集全文所有 span 的 size，取 P75 百分位作为初始阈值
    2. 若该阈值命中的 block 占比 > 50%，逐级尝试更高字号
    3. 取第一个命中数 > 2 且占比 ≤ 50% 的字号作为标题阈值
    4. 所有字号均不满足时，整文件作一个 section（section_id = "doc"）
    """
    total_pages = len(doc)
    if total_pages == 0:
        return []

    # ── 第一遍：收集所有 span size ──────────────────────────────────────
    all_sizes: List[float] = []
    for page in doc:
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    size = span.get("size", 0)
                    if size > 0:
                        all_sizes.append(size)

    # ── 第二遍：按页收集段落记录（携带 max_size） ────────────────────────
    para_records: List[dict] = []
    for page_idx, page in enumerate(doc):
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            block_text       = ""
            max_size_in_block = 0.0
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    block_text += span.get("text", "")
                    size = span.get("size", 0)
                    if size > max_size_in_block:
                        max_size_in_block = size

            block_text = block_text.strip()
            if not block_text:
                continue

            para_records.append({
                "max_size": max_size_in_block,
                "text":     block_text,
                "page":     page_idx + 1,
            })

    # ── 第三步：确定标题阈值（逐级上探） ─────────────────────────────────
    title_threshold = _find_title_threshold(all_sizes, para_records)

    # 根据最终阈值标记每个 block 的 is_title
    for rec in para_records:
        rec["is_title"] = (title_threshold > 0
                           and rec["max_size"] >= title_threshold
                           and len(rec["text"]) <= PDF_TITLE_MAX_CHARS)

    # ── 第四步：按标题切割 section ───────────────────────────────────────
    sections: List[dict] = []

    if not para_records or not any(r["is_title"] for r in para_records):
        full_text  = "\n".join(r["text"] for r in para_records)
        first_page = para_records[0]["page"] if para_records else 1
        sections.append({"section_id": "doc",
                          "page_or_sheet": str(first_page),
                          "text": full_text,
                          "section_title": ""})
        return sections

    current_title = ""
    current_page  = para_records[0]["page"]
    body_parts: List[str] = []
    counter       = 0

    def _flush():
        nonlocal counter
        body     = "\n".join(body_parts).strip()
        sec_text = (current_title + "\n" + body).strip() if current_title else body
        if sec_text:
            sections.append({"section_id":    f"s{counter}",
                              "page_or_sheet": str(current_page),
                              "text":          sec_text,
                              "section_title": ""})
            counter += 1

    for rec in para_records:
        if rec["is_title"]:
            _flush()
            current_title = rec["text"]
            current_page  = rec["page"]
            body_parts    = []
        else:
            body_parts.append(rec["text"])

    _flush()

    return sections


def _extract_pdf_images(doc, file_record: dict) -> List[dict]:
    """
    从 PDF 中提取嵌入式图片，返回 image section 列表。

    逻辑：
    1. 逐页遍历，对每页调用 page.get_images(full=True) 获取图片 xref
    2. 按 xref 去重（同一张图可能被多页引用）
    3. doc.extract_image(xref) 获取原始字节
    4. 收集所有待处理图片 → 并发调用 VLM 分类+描述（+OCR）
    5. 组装为 section dict

    section_id 格式: "img_p{page}_{idx}"，区别于文本 section
    """
    seen_xrefs: set = set()
    filename    = file_record.get("file_name", "")
    folder_code = file_record.get("folder_code")
    source_path = file_record.get("relative_path", filename)
    global_img_idx = 0  # 全局图片计数器

    # Phase 1: 收集所有待处理的图片任务
    image_tasks = []   # _process_single_image 参数
    task_meta = []     # 每个任务对应的 (page_idx, global_img_idx)

    for page_idx, page in enumerate(doc):
        try:
            images = page.get_images(full=True)
        except Exception:
            continue

        for img_info in images:
            xref = img_info[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)

            try:
                base_image = doc.extract_image(xref)
                if base_image is None:
                    continue
                raw_bytes = base_image["image"]
                img_ext   = base_image.get("ext", "png")

                png_bytes, w, h = _to_png_bytes(raw_bytes, img_ext)
                image_tasks.append({
                    "png_bytes": png_bytes,
                    "width": w,
                    "height": h,
                    "filename": filename,
                    "page": str(page_idx + 1),
                    "idx": global_img_idx,
                    "folder_code": folder_code,
                    "source_path": source_path,
                })
                task_meta.append((page_idx, global_img_idx))
                global_img_idx += 1

            except Exception:
                continue  # 单张图片失败不影响其他

    # Phase 2: 并发处理（VLM 分类+描述 → 可能追加 OCR）
    results = _process_images_batch(image_tasks)

    # Phase 3: 组装 section 列表
    sections: List[dict] = []
    for (page_idx, img_idx), text in zip(task_meta, results):
        if text is None:
            continue
        sections.append({
            "section_id":    f"img_p{page_idx + 1}_{img_idx + 1}",
            "page_or_sheet": str(page_idx + 1),
            "text":          text,
            "section_title": "",
        })

    return sections


def _extract_pdf_sections(file_record: dict, img_base_dir: str | None = None) -> List[dict]:
    """
    PDF 文本提取：返回 section 列表（不分块）。含嵌入式图片提取。

    v2 路由逻辑：
      - classify_pdf_v2 → "pymupdf" → parse_normal_pdf + _extract_pdf_images（不变）
      - classify_pdf_v2 → "sdk"     → parse_pdf_sdk（SDK 内含图片处理）
    """
    try:
        import fitz
        doc   = fitz.open(file_record["file_path"])
        route = classify_pdf_v2(doc)

        if route == "sdk":
            # SDK 路径：文本 + 图片处理已内含在 parse_pdf_sdk 中
            sections = parse_pdf_sdk(file_record, img_base_dir=img_base_dir)
        else:
            # PyMuPDF 路径：完全不变
            sections = parse_normal_pdf(doc)
            sections.extend(_extract_pdf_images(doc, file_record))

        return sections
    except Exception as e:
        print(f"  [警告] PDF 提取失败：{file_record['file_name']} — {e}")
        return []


def extract_pdf(file_record: dict) -> List[dict]:
    """
    PDF 提取对外入口。

    v2 分流策略（两分类，互补）：
      - pymupdf：纯文字 PDF（非 PPT 转 PDF，且所有页均 ≥ PDF_PAGE_MIN_CHARS 字符）
                 → parse_normal_pdf() + _extract_pdf_images()
      - sdk：需要 OCR 的 PDF（PPT 转 PDF、扫描件、或混合型——存在任一页 < 阈值）
             → parse_pdf_sdk()（SDK 内含版面检测 + OCR + 图片处理）

    异常时打印警告并返回空列表，不向外抛出。
    """
    sections = _extract_pdf_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("pdf")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# ==============================================================================
# 2-2  DOCX 提取
# ==============================================================================

# 识别"标题段落"的两种信号：
#   1. 段落 style.name 包含 "Heading" 或 "标题"（标准 Word 样式）
#   2. 段落文本匹配数字/汉字编号前缀，且字数 ≤ 80
#
# 注意：此正则匹配所有层级（1. / 1.1 / 1.1.1），层级裁剪由 parse_docx 动态决定。
_HEADING_STYLE_RE = re.compile(r'[Hh]eading|标题|\bTitle\b', re.IGNORECASE)
_HEADING_TEXT_RE  = re.compile(
    r'^(?:'
    r'\d+(?:\.\d+)*[\.、\s]'       # 1. / 1.1 / 1.1.1 等多级数字编号
    r'|[一二三四五六七八九十]+[\.、\s]'  # 一、 二、 等汉字编号
    r'|第[一二三四五六七八九十百\d]+[章节条款]'  # 第一章 / 第2节
    r')'
)


def _heading_numeric_level(text: str) -> int:
    """
    从标题文本推断数字编号层级（点分层级数）。
    "1.目的" → 1，"3.1 总经理" → 2，"6.2.2.1 熟悉…" → 4
    汉字编号/第X章 → 1
    无法识别 → 0
    """
    m = re.match(r'^(\d+(?:\.\d+)*)[\.、\s]', text)
    if m:
        return len(m.group(1).split('.'))
    if re.match(r'^(?:[一二三四五六七八九十]+[\.、\s]|第[一二三四五六七八九十百\d]+[章节条款])', text):
        return 1
    return 0


def _is_heading_para(para) -> bool:
    """判断 python-docx Paragraph 是否为标题段落（任意层级）。"""
    style_name = para.style.name if para.style else ""
    if _HEADING_STYLE_RE.search(style_name):
        return True
    text = para.text.strip()
    if not text or len(text) > DOCX_HEADING_MAX_CHARS:
        return False
    return bool(_HEADING_TEXT_RE.match(text))


def _detect_max_heading_level(element_records: list) -> int:
    """
    扫描元素记录，返回所有标题段落中最深的数字编号层级。
    Word Heading 样式的段落也纳入计算（视为层级 1）。
    """
    max_lv = 0
    for rec in element_records:
        if not rec.get("is_title"):
            continue
        lv = rec.get("heading_level", 0)
        if lv > max_lv:
            max_lv = lv
    return max_lv


def _table_to_text(table) -> str:
    """
    将 python-docx Table 转换为可读文本。
    每行用制表符分隔单元格，行间用换行符分隔。
    合并单元格处理：
      - 列合并（colspan）：同行内 _tc 相同的单元格只取一次
      - 行合并（rowspan）：不同行间 _tc 可能相同，但每行独立去重，避免漏行
    """
    row_texts: List[str] = []
    for row in table.rows:
        seen_in_row: set = set()   # 只在行内去重，避免跨行合并误删整行
        cell_texts: List[str] = []
        for cell in row.cells:
            cell_id = id(cell._tc)
            if cell_id in seen_in_row:
                continue
            seen_in_row.add(cell_id)
            cell_text = cell.text.strip()
            if cell_text:
                cell_texts.append(cell_text)
        if cell_texts:
            row_texts.append("\t".join(cell_texts))
    return "\n".join(row_texts)


def _iter_body_elements(doc_obj):
    """
    按 Word 文档的原始顺序（XML 顺序）迭代段落与表格，
    保持表格在正文中的相对位置。
    每个元素返回 ("para", Paragraph) 或 ("table", Table)。
    """
    body = doc_obj.element.body
    para_iter  = iter(doc_obj.paragraphs)
    table_iter = iter(doc_obj.tables)

    cur_para  = next(para_iter,  None)
    cur_table = next(table_iter, None)

    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p" and cur_para is not None and cur_para._element is child:
            yield ("para", cur_para)
            cur_para = next(para_iter, None)
        elif tag == "tbl" and cur_table is not None and cur_table._tbl is child:
            yield ("table", cur_table)
            cur_table = next(table_iter, None)


def parse_docx(doc_obj) -> List[dict]:
    """
    对 python-docx Document 对象按标题结构切 section，返回 section 列表。

    内容来源（按 XML 顺序遍历，保留原始位置）：
      - 段落：非空段落文本
      - 表格：转换为"单元格文本 tab 分隔"的多行文本，附加到当前 section

    动态层级切割策略：
      - 探测文档最大标题层级 max_level
      - max_level ≤ 2：用"第1层"标题切 section（现有逻辑）
      - max_level ≥ 3：用"第2层"标题切 section，
          第1层标题不触发切割，记入 section_title 字段（供对齐表审查上下文）

    无任何标题时整文件作一个 section（section_id="doc"）。
    """
    # ── 按 XML 顺序收集元素记录，标注标题层级 ──────────────────────────────────
    element_records: List[dict] = []
    for kind, elem in _iter_body_elements(doc_obj):
        if kind == "para":
            text = elem.text.strip()
            if not text:
                continue
            is_title = _is_heading_para(elem)
            lv = _heading_numeric_level(text) if is_title else 0
            # Word Heading 样式但无数字编号，视为第 1 层
            if is_title and lv == 0:
                lv = 1
            element_records.append({
                "is_title":      is_title,
                "heading_level": lv,
                "text":          text,
            })
        else:  # table
            tbl_text = _table_to_text(elem)
            if not tbl_text.strip():
                continue
            element_records.append({
                "is_title":      False,
                "heading_level": 0,
                "text":          tbl_text,
            })

    if not element_records:
        return []

    # ── 无标题：整文件作一个 section ─────────────────────────────────────────
    if not any(r["is_title"] for r in element_records):
        full_text = "\n".join(r["text"] for r in element_records)
        sections = [{"section_id": "doc", "page_or_sheet": "1",
                     "text": full_text, "section_title": ""}]
        return sections

    # ── 动态确定切割层级 ──────────────────────────────────────────────────────
    max_level = _detect_max_heading_level(element_records)
    # max_level ≤ 2：用第1层切；max_level ≥ 3：用第2层切
    cut_level = 1 if max_level <= 2 else 2

    # ── 按 cut_level 切割 section ────────────────────────────────────────────
    sections:      List[dict] = []
    current_l1     = ""   # 第1层标题（仅 cut_level=2 时用作 section_title）
    current_title  = ""   # 当前 section 的标题行
    body_parts:    List[str] = []
    counter        = 0

    def _flush():
        nonlocal counter
        body     = "\n".join(body_parts).strip()
        sec_text = (current_title + "\n" + body).strip() if current_title else body
        if sec_text:
            sections.append({
                "section_id":    f"s{counter}",
                "page_or_sheet": "1",
                "text":          sec_text,
                "section_title": current_l1 if cut_level == 2 else "",
            })
            counter += 1

    for rec in element_records:
        if not rec["is_title"]:
            body_parts.append(rec["text"])
            continue

        lv = rec["heading_level"]

        if cut_level == 1:
            # 所有标题都触发切割
            _flush()
            current_title = rec["text"]
            body_parts    = []
        else:
            # cut_level == 2
            if lv == 1:
                # 第1层：不切，更新 l1 上下文；将标题文本并入当前 body
                _flush()
                current_l1    = rec["text"]
                current_title = ""
                body_parts    = [rec["text"]]  # l1 标题保留在正文开头
            elif lv == 2:
                # 第2层：触发切割
                _flush()
                current_title = rec["text"]
                body_parts    = []
            else:
                # 第3层及以下：不切，作为正文保留
                body_parts.append(rec["text"])

    _flush()

    return sections


def _extract_docx_images(file_path: str, file_record: dict) -> List[dict]:
    """
    从 DOCX 中提取嵌入式图片（word/media/ 目录下的文件）。

    原理：.docx 是 ZIP 格式，word/media/ 下存放所有嵌入图片。
    直接解压读取，避免依赖 InlineShape（python-docx 对图片支持有限）。
    跳过 EMF/WMF 矢量格式（Pillow 不支持）。

    返回 image section 列表，section_id 格式："img_d_{idx}"
    """
    import zipfile

    filename    = file_record.get("file_name", "")
    folder_code = file_record.get("folder_code")
    source_path = file_record.get("relative_path", filename)
    _IMG_EXTS = {"png", "jpg", "jpeg", "bmp", "tiff", "gif"}

    # Phase 1: 收集图片任务
    image_tasks = []
    task_indices = []  # 对应的 idx

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            media_files = sorted([
                n for n in zf.namelist()
                if n.startswith("word/media/")
                and n.rsplit(".", 1)[-1].lower() in _IMG_EXTS
            ])

            for idx, media_path in enumerate(media_files):
                try:
                    raw_bytes = zf.read(media_path)
                    media_ext = media_path.rsplit(".", 1)[-1].lower()

                    png_bytes, w, h = _to_png_bytes(raw_bytes, media_ext)
                    image_tasks.append({
                        "png_bytes": png_bytes,
                        "width": w,
                        "height": h,
                        "filename": filename,
                        "page": "1",
                        "idx": idx,
                        "folder_code": folder_code,
                        "source_path": source_path,
                    })
                    task_indices.append(idx)
                except Exception:
                    continue
    except Exception:
        pass  # ZIP 打开失败（可能不是有效 DOCX）

    # Phase 2: 并发处理
    results = _process_images_batch(image_tasks)

    # Phase 3: 组装 section 列表
    sections: List[dict] = []
    for idx, text in zip(task_indices, results):
        if text is None:
            continue
        sections.append({
            "section_id":    f"img_d_{idx + 1}",
            "page_or_sheet": "1",
            "text":          text,
            "section_title": "",
        })

    return sections


def _extract_docx_sections(file_record: dict) -> List[dict]:
    """DOCX 文本提取：返回 section 列表（不分块）。含嵌入式图片提取。"""
    try:
        import docx
        doc = docx.Document(file_record["file_path"])
        text_sections = parse_docx(doc)
        img_sections  = _extract_docx_images(file_record["file_path"], file_record)
        return text_sections + img_sections
    except Exception as e:
        print(f"  [警告] DOCX 提取失败：{file_record['file_name']} — {e}")
        return []


def extract_docx(file_record: dict) -> List[dict]:
    """
    .docx 提取对外入口。
    使用 python-docx 解析文档结构，按标题切 section。
    异常时打印警告并返回空列表，不向外抛出。
    """
    sections = _extract_docx_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("docx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# ==============================================================================
# 2-3  DOC 提取（LibreOffice 转换）
# ==============================================================================

# LibreOffice 转换输出的临时目录
_DOC_TMP_DIR = os.path.join(os.path.dirname(__file__), '..', 'data', 'processed', '_doc_tmp')

# 懒加载警告：避免重复打印 LibreOffice 缺失提示
_soffice_warned = False


def _soffice_available() -> bool:
    """检查 soffice（LibreOffice）是否在 PATH 中可用。"""
    import shutil
    return shutil.which("soffice") is not None


def convert_doc_to_docx(doc_path: str) -> str | None:
    """
    用 LibreOffice 将 .doc 转换为 .docx，返回转换后的文件路径。
    转换失败返回 None。
    输出文件放在 _DOC_TMP_DIR 临时目录，调用方负责清理（或保留缓存）。
    """
    global _soffice_warned
    if not _soffice_available():
        if not _soffice_warned:
            print("  [警告] 未找到 soffice（LibreOffice），.doc 文件将跳过文本提取。")
            print("         安装方式：sudo apt install libreoffice-core libreoffice-writer")
            _soffice_warned = True
        return None

    import subprocess
    os.makedirs(_DOC_TMP_DIR, exist_ok=True)

    try:
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "docx",
             "--outdir", _DOC_TMP_DIR, doc_path],
            capture_output=True, text=True, timeout=SOFFICE_DOC_TIMEOUT,
        )
        if result.returncode != 0:
            print(f"  [警告] soffice 转换失败（returncode={result.returncode}）："
                  f"{os.path.basename(doc_path)}")
            return None

        # 输出文件名 = 原文件名 stem + .docx
        stem    = os.path.splitext(os.path.basename(doc_path))[0]
        out_path = os.path.join(_DOC_TMP_DIR, stem + ".docx")
        if os.path.exists(out_path):
            return out_path

        print(f"  [警告] soffice 转换后未找到输出文件：{out_path}")
        return None
    except subprocess.TimeoutExpired:
        print(f"  [警告] soffice 转换超时：{os.path.basename(doc_path)}")
        return None
    except Exception as e:
        print(f"  [警告] soffice 转换异常：{os.path.basename(doc_path)} — {e}")
        return None


def _extract_doc_sections(file_record: dict) -> List[dict]:
    """DOC 文本提取：LibreOffice 转 .docx 后提取 sections（不分块）。"""
    docx_path = convert_doc_to_docx(file_record["file_path"])
    if docx_path is None:
        return []
    tmp_record = dict(file_record)
    tmp_record["file_path"] = docx_path
    return _extract_docx_sections(tmp_record)


def extract_doc(file_record: dict) -> List[dict]:
    """
    .doc 提取对外入口。
    先用 LibreOffice 转换为 .docx，再复用 extract_docx()。
    LibreOffice 不可用或转换失败时返回空列表（路径编码标签仍保留）。
    """
    sections = _extract_doc_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("docx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# ==============================================================================
# 2-4  XLSX / XLS 提取（openpyxl / xlrd）
# ==============================================================================

def _sheet_to_text(sheet_obj, fmt: str) -> str:
    """
    将单张 sheet 转换为可读文本。

    格式：单元格间 Tab 分隔，行间换行，与 _table_to_text() 保持一致。
    空行（所有单元格均为空）跳过。

    fmt:
        "xlsx" — openpyxl Worksheet（iter_rows(values_only=True)）
        "xls"  — xlrd Sheet（row_values(i)）
    """
    row_texts: List[str] = []

    if fmt == "xlsx":
        for row in sheet_obj.iter_rows(values_only=True):
            cell_texts = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if cell_texts:
                row_texts.append("\t".join(cell_texts))
    else:  # xls
        for i in range(sheet_obj.nrows):
            cell_texts = [str(v).strip() for v in sheet_obj.row_values(i)
                          if v is not None and str(v).strip()]
            if cell_texts:
                row_texts.append("\t".join(cell_texts))

    return "\n".join(row_texts)


def parse_xlsx(workbook, fmt: str) -> List[dict]:
    """
    对 openpyxl Workbook 或 xlrd Workbook 对象，
    按 sheet 切 section，返回 section 列表。

    section 划分规则：
        每个有内容的 Sheet = 1 个 section
        section_id    = "s{i}"（i 为 sheet 的从零开始的索引）
        page_or_sheet = sheet 名称
        section_title = ""（Excel 无标题层级）

    空 sheet（text 无有效字符）在此处直接过滤。
    """
    sections: List[dict] = []

    if fmt == "xlsx":
        for i, name in enumerate(workbook.sheetnames):
            ws = workbook[name]
            text = _sheet_to_text(ws, "xlsx")
            if not count_meaningful_chars(text):
                continue  # 空 sheet，跳过
            sections.append({
                "section_id":    f"s{i}",
                "page_or_sheet": name,
                "text":          text,
                "section_title": "",
            })
    else:  # xls
        for i in range(workbook.nsheets):
            ws = workbook.sheet_by_index(i)
            text = _sheet_to_text(ws, "xls")
            if not count_meaningful_chars(text):
                continue  # 空 sheet，跳过
            sections.append({
                "section_id":    f"s{i}",
                "page_or_sheet": ws.name,
                "text":          text,
                "section_title": "",
            })

    return sections


def _extract_xlsx_sections(file_record: dict) -> List[dict]:
    """XLSX 文本提取：返回 section 列表（不分块）。"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_record["file_path"],
                                     read_only=True, data_only=True)
        return parse_xlsx(wb, "xlsx")
    except Exception as e:
        print(f"  [警告] XLSX 提取失败：{file_record['file_name']} — {e}")
        return []


def _extract_xls_sections(file_record: dict) -> List[dict]:
    """XLS 文本提取：返回 section 列表（不分块）。"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_record["file_path"])
        return parse_xlsx(wb, "xls")
    except Exception as e:
        print(f"  [警告] XLS 提取失败：{file_record['file_name']} — {e}")
        return []


def extract_xlsx(file_record: dict) -> List[dict]:
    """
    .xlsx 提取对外入口。
    使用 openpyxl 逐 sheet 读取，每 sheet 作一个 section。
    异常时打印警告并返回空列表，不向外抛出。
    """
    sections = _extract_xlsx_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("xlsx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


def extract_xls(file_record: dict) -> List[dict]:
    """
    .xls 提取对外入口。
    使用 xlrd 逐 sheet 读取，每 sheet 作一个 section。
    异常时打印警告并返回空列表，不向外抛出。
    """
    sections = _extract_xls_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("xlsx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# ==============================================================================
# 2-5  PPTX / PPT 提取（python-pptx / LibreOffice 转换）
# ==============================================================================

_PPT_TMP_DIR        = os.path.join(os.path.dirname(__file__), '..', 'data', 'processed', '_ppt_tmp')
_soffice_ppt_warned = False


def _slide_to_text(slide) -> str:
    """
    提取单张 slide 的全部文本。

    处理规则：
    - 文本框 / placeholder（has_text_frame=True）：段落间 \\n 连接，shape 间 \\n\\n 分隔
    - 分组（shape_type == 6）：递归遍历内部子 shape，子 shape 规则同上
    - 表格（shape_type == 19）：cell 间 \\t，行间 \\n，与 _table_to_text() 格式一致
    - 图片（shape_type == 13）等无文本 shape：静默跳过

    shape 提取顺序按 slide 内 XML 自然顺序（slide.shapes）。
    """
    parts: List[str] = []

    def _collect(shapes_iter) -> None:
        for shape in shapes_iter:
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if paras:
                    parts.append("\n".join(paras))
            elif shape.shape_type == 6:    # GROUP — 递归展开
                try:
                    _collect(shape.shapes)
                except Exception:
                    pass                   # 无法迭代的 GROUP 静默跳过
            elif shape.shape_type == 19:   # MSO_SHAPE_TYPE.TABLE == 19
                row_texts: List[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        row_texts.append("\t".join(cells))
                if row_texts:
                    parts.append("\n".join(row_texts))
            # shape_type == 13（PICTURE）及其他无文本 shape：静默跳过

    _collect(slide.shapes)

    return "\n\n".join(parts)


def parse_pptx(prs) -> List[dict]:
    """
    对 python-pptx Presentation 对象，按 slide 切 section，返回 section 列表。

    每张非空 slide = 1 个 section：
        section_id    = f"s{i}"（slide 索引，从 0 计数）
        page_or_sheet = str(i + 1)（物理页码，从 1 计数）
        section_title = ""（PPT 无层级标题概念）

    空 slide（有效字符为 0）提前跳过。
    """
    sections: List[dict] = []

    for i, slide in enumerate(prs.slides):
        text = _slide_to_text(slide)
        if not count_meaningful_chars(text):
            continue  # 空 slide（纯图片页等），跳过
        sections.append({
            "section_id":    f"s{i}",
            "page_or_sheet": str(i + 1),
            "text":          text,
            "section_title": "",
        })

    return sections


def _extract_pptx_images(prs, file_record: dict) -> List[dict]:
    """
    从 PPTX 中提取图片 shape（shape_type==13 即 PICTURE）。

    逐 slide 遍历，对每个 PICTURE shape 读取 shape.image.blob。
    使用 SHA256(blob) 去重（同一张图多 slide 复用时只处理一次）。
    递归处理 GROUP（shape_type==6）内嵌套的图片。

    返回 image section 列表，section_id 格式："img_s{slide}_{idx}"
    """
    import hashlib

    filename    = file_record.get("file_name", "")
    folder_code = file_record.get("folder_code")
    source_path = file_record.get("relative_path", filename)
    seen_hashes: set = set()
    global_counter = [0]  # 用列表便于闭包修改

    # Phase 1: 收集图片任务
    image_tasks = []
    task_meta = []  # (slide_idx, counter_value)

    def _collect_pictures(shapes_iter, slide_idx: int):
        for shape in shapes_iter:
            if shape.shape_type == 13:  # PICTURE
                try:
                    blob = shape.image.blob
                    blob_hash = hashlib.sha256(blob).hexdigest()
                    if blob_hash in seen_hashes:
                        continue
                    seen_hashes.add(blob_hash)

                    png_bytes, w, h = _to_png_bytes(blob)
                    image_tasks.append({
                        "png_bytes": png_bytes,
                        "width": w,
                        "height": h,
                        "filename": filename,
                        "page": str(slide_idx + 1),
                        "idx": global_counter[0],
                        "folder_code": folder_code,
                        "source_path": source_path,
                    })
                    global_counter[0] += 1
                    task_meta.append((slide_idx, global_counter[0]))
                except Exception as e:
                    print(f"[WARN] 跳过 PPTX 图片（slide {slide_idx + 1}）：{e}")
                    continue
            elif shape.shape_type == 6:  # GROUP — 递归展开
                try:
                    _collect_pictures(shape.shapes, slide_idx)
                except Exception:
                    pass

    for i, slide in enumerate(prs.slides):
        _collect_pictures(slide.shapes, i)

    # Phase 2: 并发处理
    results = _process_images_batch(image_tasks)

    # Phase 3: 组装 section 列表
    sections: List[dict] = []
    for (slide_idx, counter_val), text in zip(task_meta, results):
        if text is None:
            continue
        sections.append({
            "section_id":    f"img_s{slide_idx + 1}_{counter_val}",
            "page_or_sheet": str(slide_idx + 1),
            "text":          text,
            "section_title": "",
        })

    return sections


def _extract_pptx_sections(file_record: dict) -> List[dict]:
    """PPTX 文本提取：返回 section 列表（不分块）。含图片提取。"""
    try:
        from pptx import Presentation
        prs = Presentation(file_record["file_path"])
        text_sections = parse_pptx(prs)
        img_sections  = _extract_pptx_images(prs, file_record)
        return text_sections + img_sections
    except Exception as e:
        print(f"  [警告] PPTX 提取失败：{file_record['file_name']} — {e}")
        return []


def _extract_ppt_sections(file_record: dict) -> List[dict]:
    """PPT 文本提取：LibreOffice 转 .pptx 后提取 sections（不分块）。"""
    pptx_path = convert_ppt_to_pptx(file_record["file_path"])
    if pptx_path is None:
        return []
    tmp_record = dict(file_record)
    tmp_record["file_path"] = pptx_path
    return _extract_pptx_sections(tmp_record)


def extract_pptx(file_record: dict) -> List[dict]:
    """
    .pptx 提取对外入口。
    使用 python-pptx 逐 slide 读取，每 slide 作一个 section。
    异常时打印警告并返回空列表，不向外抛出。
    """
    sections = _extract_pptx_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("pptx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


def convert_ppt_to_pptx(ppt_path: str) -> str | None:
    """
    用 LibreOffice 将 .ppt 转换为 .pptx，返回转换后的文件路径。
    转换失败返回 None。
    输出文件放在 _PPT_TMP_DIR 临时目录，调用方负责清理（或保留缓存）。
    """
    global _soffice_ppt_warned
    if not _soffice_available():
        if not _soffice_ppt_warned:
            print("  [警告] 未找到 soffice（LibreOffice），.ppt 文件将跳过文本提取。")
            print("         安装方式：sudo apt install libreoffice-core libreoffice-impress")
            _soffice_ppt_warned = True
        return None

    import subprocess
    os.makedirs(_PPT_TMP_DIR, exist_ok=True)

    try:
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pptx",
             "--outdir", _PPT_TMP_DIR, ppt_path],
            capture_output=True, text=True, timeout=SOFFICE_PPT_TIMEOUT,
        )
        if result.returncode != 0:
            print(f"  [警告] soffice 转换失败（returncode={result.returncode}）："
                  f"{os.path.basename(ppt_path)}")
            return None

        stem     = os.path.splitext(os.path.basename(ppt_path))[0]
        out_path = os.path.join(_PPT_TMP_DIR, stem + ".pptx")
        if os.path.exists(out_path):
            return out_path

        print(f"  [警告] soffice 转换后未找到输出文件：{out_path}")
        return None
    except subprocess.TimeoutExpired:
        print(f"  [警告] soffice 转换超时：{os.path.basename(ppt_path)}")
        return None
    except Exception as e:
        print(f"  [警告] soffice 转换异常：{os.path.basename(ppt_path)} — {e}")
        return None


def extract_ppt(file_record: dict) -> List[dict]:
    """
    .ppt 提取对外入口。
    先用 LibreOffice 转换为 .pptx，再复用 extract_pptx()。
    LibreOffice 不可用或转换失败时返回空列表（路径编码标签仍保留）。
    """
    sections = _extract_ppt_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("pptx")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# ==============================================================================
# 2-6  JPG / PNG 提取（GLM-OCR）
# ==============================================================================

def _extract_image_sections(file_record: dict) -> List[dict]:
    """
    独立图片提取（JPG/PNG）— VLM 分类+描述优先，文档/表格使用 SDK OCR。

    改进点（相比原版统一走 GLM-OCR）：
    - 原版：所有图片统一送 GLM-OCR → 照片/LOGO 返回乱码
    - 新版：VLM 分类+描述 → 文档/表格类使用 SDK 版面检测 + 结构化 OCR
    - VLM 不可用时回退到 SDK OCR（利用版面检测能力）
    """
    try:
        from PIL import Image
        import io

        img_path    = file_record["file_path"]
        filename    = file_record.get("file_name", "")
        folder_code = file_record.get("folder_code")
        source_path = file_record.get("relative_path", filename)

        with Image.open(img_path) as img:
            w, h = img.size
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            png_bytes = buf.getvalue()

        # 公共管线：过滤 → VLM 分类+描述 → 组装
        text = _process_single_image(
            png_bytes, w, h,
            filename=filename,
            page="1",
            idx=0,
            folder_code=folder_code,
            source_path=source_path,
        )

        if text is None:
            # VLM 不可用或被过滤 → 回退到 SDK OCR（利用版面检测能力）
            if _filter_image(png_bytes, w, h):
                text = ocr_image_with_sdk(png_bytes, img_type="独立图片")
            else:
                return []

        return [{
            "section_id":    "s0",
            "page_or_sheet": "1",
            "text":          text,
            "section_title": "",
        }]
    except Exception as e:
        print(f"  [警告] 图片提取失败：{file_record['file_name']} — {e}")
        return []


def extract_image(file_record: dict) -> List[dict]:
    """
    JPG / JPEG / PNG 提取对外入口。

    处理流程：
      1. VLM 分类+描述 → 文档/表格类使用 SDK 版面检测 + 结构化 OCR
      2. VLM 不可用时回退到 SDK OCR

    SDK 初始化采用单例模式，首次加载约 4 秒，后续复用。
    """
    sections = _extract_image_sections(file_record)
    if not sections:
        return []
    _max, _min = chunk_params("image")
    return make_chunks_from_sections(sections, file_record,
                                     max_size=_max, min_size=_min)


# 方便按扩展名路由的别名
extract_jpg  = extract_image
extract_jpeg = extract_image
extract_png  = extract_image


# ==============================================================================
# 统一 section 提取入口（供 align_evidence.py 阶段 2a 使用）
# ==============================================================================

_SECTION_EXTRACTOR_MAP = {
    ".pdf":  _extract_pdf_sections,
    ".docx": _extract_docx_sections,
    ".doc":  _extract_doc_sections,
    ".xlsx": _extract_xlsx_sections,
    ".xls":  _extract_xls_sections,
    ".pptx": _extract_pptx_sections,
    ".ppt":  _extract_ppt_sections,
    ".jpg":  _extract_image_sections,
    ".jpeg": _extract_image_sections,
    ".png":  _extract_image_sections,
}


def extract_sections(file_record: dict, img_base_dir: str | None = None) -> list[dict] | None:
    """
    纯文本提取：返回 section 列表（不分块）。
    按文件扩展名分发到对应的 _extract_xxx_sections 函数。

    返回值：
        list[dict] — section 列表（可能为空列表 [] 表示提取失败或无内容）
        None       — 不支持的文件格式
    """
    ext = file_record.get("extension", "")
    extractor = _SECTION_EXTRACTOR_MAP.get(ext)
    if extractor is None:
        return None
    if ext == ".pdf":
        return extractor(file_record, img_base_dir=img_base_dir)
    return extractor(file_record)


# ==============================================================================
# Phase 2：Embedding 与上下文辅助函数
# ==============================================================================

def get_text_for_embedding(chunk: dict) -> str:
    """
    根据 chunk 类型选择用于 Embedding 的文本。

    规则：
    - 表格 chunk 且有 table_summary：使用纯摘要（语义密度高）
    - 其他情况：使用 text 字段

    用于 align_evidence.py 阶段 3 的 embed_chunks()
    """
    if chunk.get("is_table") and chunk.get("table_summary"):
        return chunk["table_summary"]
    return chunk.get("text", "")


def get_chunk_context(
    chunk: dict,
    parents: dict,
    max_parent_len: int = 2000,
    context_chars: int = 300,
) -> str:
    """
    获取 chunk 的上下文文本（用于 Reranker/LLM）。

    参数:
        chunk: chunk 记录
        parents: {parent_id: parent_text} 字典
        max_parent_len: parent_text 长度阈值
        context_chars: 超过阈值时的上下文窗口

    返回:
        上下文文本（Markdown 格式）
        - parent_text ≤ max_parent_len：返回完整 parent
        - parent_text > max_parent_len：返回 chunk 前后各 context_chars 字
    """
    parent_id = chunk.get("parent_id", "")
    parent_text = parents.get(parent_id, "")

    if not parent_text:
        return chunk.get("text", "")

    # parent 较短，直接返回完整内容
    if len(parent_text) <= max_parent_len:
        return parent_text

    # parent 过长，定位 chunk 并取前后文
    chunk_text = chunk.get("text", "")

    # 对于表格 chunk，用 table_markdown 定位更准确
    if chunk.get("is_table"):
        chunk_text = chunk.get("table_markdown", chunk_text)

    pos = parent_text.find(chunk_text)
    if pos == -1:
        # 尝试模糊匹配（前几行）
        lines = chunk_text.split("\n")
        first_lines = "\n".join(lines[:min(3, len(lines))])
        if first_lines:
            pos = parent_text.find(first_lines)

    if pos == -1:
        return chunk.get("text", "")

    start = max(0, pos - context_chars)
    end = min(len(parent_text), pos + len(chunk_text) + context_chars)

    return parent_text[start:end]

