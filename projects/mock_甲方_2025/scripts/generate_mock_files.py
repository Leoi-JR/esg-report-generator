"""
模拟甲方文件生成脚本
生成各类型 ESG 资料文件（DOCX/XLSX/PDF/PPTX/PNG），放置到对应文件夹。
包含 2 个故意放错位置的文件，模拟甲方操作失误。

运行方式：
    python3 projects/mock_甲方_2025/scripts/generate_mock_files.py
"""

import sys
import os
from pathlib import Path

# 项目根目录
ROOT = Path(__file__).resolve().parents[2]  # projects/mock_甲方_2025 -> projects -> repo root
# 不对：scripts -> mock_甲方_2025 -> projects -> repo root
ROOT = Path(__file__).resolve().parents[3]
BASE_DIR = ROOT / "projects" / "mock_甲方_2025" / "raw" / "整理后资料"

# ── 各编码对应的文件夹路径 ──────────────────────────────────────────────────
FOLDER_MAP = {
    "GA1": BASE_DIR / "G-公司治理/GA-治理体系/GA1 公司治理架构",
    "GA2": BASE_DIR / "G-公司治理/GA-治理体系/GA2 董事会多元化",
    "GA3": BASE_DIR / "G-公司治理/GA-治理体系/GA3 股东权益保护",
    "GC1": BASE_DIR / "G-公司治理/GC-合规建设/GC1 反腐败管理",
    "GC2": BASE_DIR / "G-公司治理/GC-合规建设/GC2 合规培训",
    "EA1": BASE_DIR / "E-环境保护/EA-应对气候变化/EA1 碳排放目标",
    "EA2": BASE_DIR / "E-环境保护/EA-应对气候变化/EA2 温室气体排放",
    "EA3": BASE_DIR / "E-环境保护/EA-应对气候变化/EA3 减碳措施",
    "EC1": BASE_DIR / "E-环境保护/EC-资源效率/EC1 能源消耗",
    "EC2": BASE_DIR / "E-环境保护/EC-资源效率/EC2 可再生能源",
    "ED1": BASE_DIR / "E-环境保护/ED-环境合规/ED1 废弃物管理",
    "DA1": BASE_DIR / "D-产业价值/DA-创新驱动/DA1 研发投入",
    "DA2": BASE_DIR / "D-产业价值/DA-创新驱动/DA2 知识产权",
    "DC1": BASE_DIR / "D-产业价值/DC-产品质量/DC1 质量管理",
    "DD1": BASE_DIR / "D-产业价值/DD-数据安全/DD1 信息安全",
    "SA1": BASE_DIR / "S-人权与社会/SA-社会贡献/SA1 公益慈善",
    "SB1": BASE_DIR / "S-人权与社会/SB-社区关系/SB1 社区沟通",
    "SC1": BASE_DIR / "S-人权与社会/SC-劳工权益/SC1 员工结构",
    "SC2": BASE_DIR / "S-人权与社会/SC-劳工权益/SC2 薪酬福利",
    "SD1": BASE_DIR / "S-人权与社会/SD-职业健康安全/SD1 安全生产",
}

# ── DOCX 生成 ──────────────────────────────────────────────────────────────

def make_docx(path: Path, title: str, sections: list[tuple[str, str]]):
    """生成包含标题+多节正文的 Word 文档"""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 文档标题
    h = doc.add_heading(title, level=1)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()  # 空行

    for section_title, content in sections:
        doc.add_heading(section_title, level=2)
        p = doc.add_paragraph(content)
        p.paragraph_format.space_after = Pt(6)

    doc.save(str(path))
    print(f"  [DOCX] {path.relative_to(BASE_DIR)}")


# ── 各编码的 DOCX 内容定义 ─────────────────────────────────────────────────

DOCX_CONTENTS = {
    "GA1": (
        "公司治理架构说明书",
        [
            ("一、股东会职能与会议机制",
             "本公司股东会为最高权力机构，负责审议批准年度报告、利润分配方案及重大投资决策。"
             "2024年度共召开股东会3次，其中年度股东会1次，临时股东会2次，出席率均超过90%。"
             "所有股东会决议均经合法程序表决通过，充分保障各类股东的知情权和参与权。"),
            ("二、董事会构成与运作",
             "公司董事会由9名董事组成，其中独立董事3名，占比33.3%，符合监管要求。"
             "董事会下设审计委员会、薪酬与考核委员会、提名委员会三个专门委员会。"
             "2024年度召开董事会会议12次，审议议题52项，审计委员会召开会议4次。"),
            ("三、监事会职责履行",
             "监事会由5名监事组成，包括股东代表监事3名和职工代表监事2名。"
             "2024年度召开监事会会议6次，对公司财务状况、高管履职情况实施有效监督，"
             "未发现公司高级管理人员存在违反法律法规及公司章程的情况。"),
        ]
    ),
    "GA2": (
        "董事会多元化政策与执行报告",
        [
            ("一、多元化政策概述",
             "公司董事会高度重视成员多元化建设，已于2022年制定并发布《董事会多元化政策》，"
             "从性别、专业背景、行业经验、国际化视野等多个维度推进多元化目标。"),
            ("二、性别多元化现状",
             "截至2024年12月31日，公司董事会共有女性董事2名，占全体董事的22.2%，"
             "较上年提升5.6个百分点。公司承诺到2026年将女性董事比例提升至30%以上。"),
            ("三、专业背景分布",
             "现任董事中，具有财务会计背景3人、法律合规背景2人、行业技术背景2人、"
             "战略管理背景2人，专业互补性强，有利于多角度审视重大决策。"),
        ]
    ),
    "GC1": (
        "反腐败与廉洁合规管理制度",
        [
            ("一、制度体系建设",
             "公司已建立以《廉洁从业行为准则》为核心的反腐败制度体系，涵盖商业贿赂预防、"
             "利益冲突管理、礼品与款待规范、举报与调查程序四大模块。所有制度每两年更新一次。"),
            ("二、举报渠道与案件处理",
             "公司设立独立举报热线（400-XXX-XXXX）及匿名举报邮箱，确保举报人身份保密。"
             "2024年度收到举报线索8件，经核查：2件属实并依规处理，4件不属实，2件转至相关部门。"),
            ("三、年度合规培训",
             "2024年对全体员工开展反腐败专题培训，覆盖率100%，考核通过率98.7%。"
             "对重点岗位（采购、财务、销售）员工额外开展深度合规培训。"),
        ]
    ),
    "EA1": (
        "碳中和目标与气候承诺声明",
        [
            ("一、碳排放绝对目标",
             "公司承诺：以2020年为基准年，到2030年实现范围一和范围二碳排放总量较基准年"
             "降低50%；到2050年实现全价值链净零碳排放，与《巴黎协定》1.5°C路径保持一致。"),
            ("二、近期行动计划（2024-2026）",
             "阶段目标：2025年较2020年减排30%。主要措施包括：工厂屋顶光伏装机扩容至50MW、"
             "高能耗设备电气化改造（覆盖80%生产线）、绿电采购比例提升至40%。"),
            ("三、第三方核查",
             "公司聘请SGS对2024年度碳排放数据进行独立核查，核查范围涵盖国内全部8家生产基地，"
             "核查声明已在公司官网及港交所ESG报告中公开披露。"),
        ]
    ),
    "EA2": (
        "2024年度温室气体排放数据报告",
        [
            ("一、排放边界与方法论",
             "本报告采用GHG Protocol企业标准，运营控制权法确定排放边界，"
             "涵盖国内8家生产基地及3家研发中心。全球变暖潜能值（GWP）采用IPCC AR6数值。"),
            ("二、范围一排放",
             "2024年范围一直接排放总量：12,450吨CO₂e，主要来源为天然气燃烧（68%）"
             "和公司自有车辆（32%）。较2023年下降8.3%，较基准年2020年下降22%。"),
            ("三、范围二排放",
             "基于市场法：18,320吨CO₂e（含绿电采购抵消后净值8,120吨）。"
             "基于位置法：21,650吨CO₂e。绿电采购量：1.2亿度，占总用电量35%。"),
        ]
    ),
    "EC1": (
        "能源消耗管理年度报告",
        [
            ("一、能源结构概况",
             "2024年公司综合能耗折标准煤：8.2万吨，同比下降5.1%。能源种类包括："
             "电力（占比62%）、天然气（占比28%）、柴油（占比7%）、其他（占比3%）。"),
            ("二、能源强度指标",
             "万元产值综合能耗：0.086吨标煤/万元，较上年下降6.5%，优于行业平均水平18%。"
             "单位产品能耗：主导产品A类较上年下降4.2%，B类较上年下降7.8%。"),
            ("三、节能项目实施",
             "2024年实施重点节能项目12项，合计节能量6,800吨标煤，节约能源费用约1,450万元。"
             "主要项目：余热回收利用系统（节能2,100吨标煤）、变频改造（节能1,800吨标煤）。"),
        ]
    ),
    "DA1": (
        "研发投入与创新成果报告",
        [
            ("一、研发投入总量",
             "2024年公司研发投入合计：15.8亿元，占营业收入比例：6.3%，"
             "同比增长18.2%，连续5年保持两位数增长。研发人员数量：2,840人，"
             "占员工总数的23.1%，其中硕士及以上学历占比72%。"),
            ("二、研发方向布局",
             "核心技术研发（60%）：聚焦下一代产品平台技术突破；"
             "应用技术研发（30%）：面向客户定制化需求的解决方案开发；"
             "前沿探索研发（10%）：量子计算、AI赋能等前沿技术储备。"),
            ("三、创新成果",
             "2024年获授权专利382件，其中发明专利286件，PCT国际专利45件。"
             "主导或参与制定国家/行业标准12项，荣获国家级科技奖励2项。"),
        ]
    ),
    "SC1": (
        "员工结构与人才发展报告",
        [
            ("一、员工总体构成",
             "截至2024年12月31日，公司全球员工总数：12,286人，"
             "其中中国大陆员工：9,842人（占80.1%），海外员工：2,444人（占19.9%）。"
             "正式员工：11,950人，劳务派遣：336人。"),
            ("二、性别与年龄分布",
             "女性员工占比：34.2%，管理层女性占比：28.6%。"
             "年龄分布：30岁以下占22%，30-40岁占38%，40-50岁占28%，50岁以上占12%。"
             "平均年龄：38.4岁，平均司龄：7.2年。"),
            ("三、员工流动情况",
             "2024年员工主动离职率：8.3%，低于行业平均水平（11.2%）。"
             "新入职员工：1,245人，其中应届生占比：32%。"),
        ]
    ),
    "SD1": (
        "职业健康与安全生产年度报告",
        [
            ("一、安全生产基本情况",
             "2024年公司实现连续1,825天（5年）无重大生产安全事故。"
             "工伤事故发生率（TRIR）：0.42，较上年下降15%，优于行业基准（0.85）51%。"
             "职业病确诊人数：0人。"),
            ("二、安全管理体系",
             "公司已通过ISO 45001:2018职业健康安全管理体系认证，覆盖全部生产基地。"
             "建立安全领导力积分制度，全体管理人员须完成不少于20小时/年的安全领导力培训。"),
            ("三、应急演练",
             "2024年组织全厂综合应急演练4次，专项应急演练（消防/化学品泄漏/高空作业）共26次，"
             "参与人次超过2.8万。应急预案更新率：100%，应急物资储备率：100%。"),
        ]
    ),
}

# 其余编码使用通用模板
GENERIC_DOCX = {
    "GA3": ("股东权益保护机制说明",
            [("股东大会投票机制",
              "公司建立累积投票制，保障中小股东提名董事的权利。所有股东大会议题提前15日通过"
              "官网和交易所平台公告，确保股东充分知情。2024年累计保护中小股东权益事项8项。"),
             ("分红政策",
              "公司《利润分配管理办法》规定，在满足正常经营所需资金的前提下，"
              "每年以现金方式分配的利润不低于当年实现的可分配利润的30%。"
              "2024年度现金分红总额：3.2亿元，每股分红：0.32元。")]),
    "GC2": ("合规培训体系建设报告",
            [("培训体系架构",
              "公司建立三级合规培训体系：新员工入职合规培训（必修16学时）、"
              "全员年度合规更新培训（必修4学时）、高风险岗位专项培训（必修8学时）。"),
             ("2024年培训数据",
              "全年共开展合规培训课程47期，累计参训11,820人次，覆盖率99.4%，"
              "考核通过率97.8%。在线学习平台新增合规课程18门，累计学习时长45,000小时。")]),
    "EA3": ("减碳技术措施实施方案",
            [("能效提升措施",
              "2024年完成老旧高耗能设备替换工程，涉及压缩机、泵组、风机等共计346台套，"
              "年节电量估算达1,380万度，折合减排约6,900吨CO₂。"),
             ("可再生能源替代",
              "新增屋顶分布式光伏装机容量12MW，全年发电量1,150万度，"
              "自发自用率88%，余电上网。光伏系统全生命周期减排约28万吨CO₂。")]),
    "EC2": ("可再生能源使用报告",
            [("绿电采购情况",
              "2024年通过绿电交易市场采购绿证1.2亿度，等效替代传统电力，"
              "减少碳排放约6.6万吨CO₂。绿证来源：风电（65%）、光伏（35%）。"),
             ("自建可再生能源",
              "公司在5个生产基地部署屋顶光伏，总装机38MW，2024年发电量3,200万度。"
              "在建项目：厂区地面光伏20MW，预计2025年Q2并网。")]),
    "ED1": ("废弃物管理与环保合规报告",
            [("废弃物产生情况",
              "2024年公司一般固体废弃物产生量：4,820吨，危险废弃物产生量：186吨。"
              "一般固废综合利用率：94.2%，危废合规处置率：100%。"),
             ("合规情况",
              "2024年未发生环保行政处罚事件，例行环保检查全部合格。"
              "完成排污许可证年度执行报告填报，主要污染物排放均达标。")]),
    "DA2": ("知识产权管理报告",
            [("专利资产现状",
              "截至2024年底，公司有效专利总量：2,186件，其中发明专利1,024件（占比46.8%）、"
              "实用新型843件、外观设计319件。PCT国际专利申请累计：312件，覆盖26个国家/地区。"),
             ("知识产权保护行动",
              "2024年发起专利维权诉讼3件，获赔金额合计820万元。完成商标注册115件，"
              "软件著作权登记42件。对新招聘研发人员全员开展知识产权保密培训。")]),
    "DC1": ("产品质量管理体系报告",
            [("质量管理认证",
              "公司已通过ISO 9001:2015质量管理体系认证、IATF 16949汽车质量管理认证，"
              "覆盖全部生产基地。2024年监督审核零不符合项，一次性通过率100%。"),
             ("质量绩效指标",
              "2024年产品出厂合格率：99.97%，客户投诉率较上年下降23%。"
              "实施全面质量管控（TQC）专项项目12个，累计节约质量损失成本780万元。")]),
    "DD1": ("信息安全与数据保护管理报告",
            [("安全管理框架",
              "公司依据ISO/IEC 27001:2022建立信息安全管理体系，2024年完成年度认证审核。"
              "设立信息安全委员会，由CTO担任主席，每季度召开专题会议审查安全态势。"),
             ("安全事件与应对",
              "2024年共发生低级别信息安全事件7起，无重大数据泄露事件。"
              "开展全员网络安全意识培训，覆盖率100%，钓鱼邮件点击率从18%降至4.2%。")]),
    "SA1": ("社会公益与慈善贡献报告",
            [("公益资金投入",
              "2024年公司及旗下基金会累计公益捐款：3,680万元，受益群体超过12万人次。"
              "重点领域：乡村教育（45%）、科技人才培养（30%）、灾害救援（15%）、其他（10%）。"),
             ("员工志愿者服务",
              "2024年员工志愿服务累计时长：18,400小时，参与员工：3,240人次，"
              "组织志愿活动42场，涵盖助老敬老、环保清洁、支教助学等主题。")]),
    "SB1": ("社区沟通与利益相关方参与报告",
            [("社区沟通机制",
              "公司建立四级社区沟通机制：日常信访受理（全年无休）、月度社区座谈会、"
              "季度开放日、年度利益相关方大会。2024年社区信访受理128件，100%结案回复。"),
             ("重要利益相关方参与",
              "2024年开展利益相关方重要性评估，有效回收问卷2,450份，涵盖投资者、"
              "客户、供应商、员工、社区居民、政府机构六大类别。")]),
    "SC2": ("薪酬福利体系与员工关怀报告",
            [("薪酬竞争力",
              "公司薪酬定位为行业75分位，每年开展市场薪酬对标调研，确保各岗位薪酬具有竞争力。"
              "2024年人均薪酬较上年增长8.3%，高于CPI增幅。"),
             ("员工福利项目",
              "公司提供补充医疗保险（含家属）、员工持股计划、带薪年假（最高20天）、"
              "弹性工作制、子女教育补贴等特色福利。员工满意度调查综合得分：82.4分（百分制）。")]),
}

# 合并所有DOCX内容定义
ALL_DOCX = {**DOCX_CONTENTS, **GENERIC_DOCX}


# ── XLSX 生成 ──────────────────────────────────────────────────────────────

def make_xlsx(path: Path, sheet_configs: list[tuple[str, list, list[list]]]):
    """
    生成多 Sheet 的 Excel 文件。
    sheet_configs: [(sheet_name, [col_headers], [[row_data], ...]), ...]
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # 删除默认空 sheet

    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=11)
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    for sheet_name, headers, rows in sheet_configs:
        ws = wb.create_sheet(title=sheet_name)
        # 写表头
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = thin_border
        # 写数据行
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = thin_border
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                if row_idx % 2 == 0:
                    cell.fill = PatternFill(start_color="EBF3FB", end_color="EBF3FB", fill_type="solid")
        # 自动列宽
        for col_idx, header in enumerate(headers, 1):
            col_letter = get_column_letter(col_idx)
            max_len = max(len(str(header)), max((len(str(r[col_idx-1])) for r in rows), default=0))
            ws.column_dimensions[col_letter].width = min(max_len * 1.5 + 4, 40)
        ws.row_dimensions[1].height = 24

    wb.save(str(path))
    print(f"  [XLSX] {path.relative_to(BASE_DIR)}")


XLSX_CONTENTS = {
    "EA2": (
        "EA2_温室气体排放数据表.xlsx",
        [
            ("范围一二排放汇总",
             ["排放类别", "排放源", "2022年(tCO₂e)", "2023年(tCO₂e)", "2024年(tCO₂e)", "同比变化"],
             [
                 ["范围一", "天然气燃烧", 9820, 9250, 8470, "-8.4%"],
                 ["范围一", "公司自有车辆", 4320, 4110, 3980, "-3.2%"],
                 ["范围一", "小计", 14140, 13360, 12450, "-8.3%"],
                 ["范围二(市场法)", "外购电力", 22180, 20640, 18320, "-11.2%"],
                 ["范围二(位置法)", "外购电力", 24650, 23210, 21650, "-6.7%"],
                 ["合计(市场法)", "", 36320, 34000, 30770, "-9.5%"],
             ]),
            ("分基地排放",
             ["基地名称", "所在省份", "范围一(tCO₂e)", "范围二-市场法(tCO₂e)", "合计(tCO₂e)"],
             [
                 ["上海基地", "上海", 2810, 4120, 6930],
                 ["苏州基地", "江苏", 3240, 5680, 8920],
                 ["深圳基地", "广东", 2190, 3840, 6030],
                 ["成都基地", "四川", 1560, 2890, 4450],
                 ["武汉基地", "湖北", 980, 1420, 2400],
                 ["其他基地(3个)", "各省", 1670, 370, 2040],
             ]),
        ]
    ),
    "EC1": (
        "EC1_能源消耗明细数据.xlsx",
        [
            ("年度能耗汇总",
             ["能源类型", "单位", "2022年", "2023年", "2024年", "折标煤(2024)"],
             [
                 ["外购电力", "万度", 18420, 17680, 16950, "20831吨标煤"],
                 ["天然气", "万立方米", 682, 651, 623, "8106吨标煤"],
                 ["柴油", "吨", 1240, 1180, 1095, "1603吨标煤"],
                 ["汽油", "吨", 328, 312, 289, "423吨标煤"],
                 ["综合能耗合计", "吨标煤", 87600, 86300, 82000, "82000吨标煤"],
             ]),
            ("节能项目清单",
             ["项目名称", "实施基地", "完成时间", "节能量(吨标煤/年)", "投资额(万元)", "回收期(年)"],
             [
                 ["余热回收利用系统", "苏州基地", "2024-03", 2100, 480, 2.3],
                 ["高效变频改造", "上海基地", "2024-06", 1800, 320, 1.8],
                 ["LED照明全替换", "深圳基地", "2024-08", 680, 95, 1.4],
                 ["空压机节能改造", "成都基地", "2024-10", 920, 180, 2.0],
                 ["屋顶隔热改造", "武汉基地", "2024-11", 450, 68, 1.5],
                 ["其他节能项目(7项)", "各基地", "全年", 850, 210, 2.5],
             ]),
        ]
    ),
    "SC1": (
        "SC1_员工结构数据表.xlsx",
        [
            ("员工总览",
             ["维度", "类别", "人数", "占比", "较上年变化"],
             [
                 ["性别", "男性", 8084, "65.8%", "+124人"],
                 ["性别", "女性", 4202, "34.2%", "+87人"],
                 ["年龄", "30岁以下", 2703, "22.0%", "+215人"],
                 ["年龄", "30-40岁", 4669, "38.0%", "+48人"],
                 ["年龄", "40-50岁", 3440, "28.0%", "-32人"],
                 ["年龄", "50岁以上", 1474, "12.0%", "-20人"],
                 ["学历", "博士及以上", 492, "4.0%", "+38人"],
                 ["学历", "硕士", 2457, "20.0%", "+156人"],
                 ["学历", "本科", 6143, "50.0%", "+102人"],
                 ["学历", "大专及以下", 3194, "26.0%", "-85人"],
             ]),
            ("各部门人员分布",
             ["部门", "人数", "女性占比", "平均年龄", "平均司龄"],
             [
                 ["研发中心", 2840, "28.4%", 35.2, 5.8],
                 ["生产制造", 4680, "31.2%", 40.1, 9.2],
                 ["销售市场", 1920, "42.6%", 36.8, 6.4],
                 ["财务法务", 680, "58.8%", 38.4, 8.1],
                 ["供应链", 1240, "22.4%", 39.2, 7.6],
                 ["职能支撑", 926, "48.2%", 37.6, 6.9],
             ]),
        ]
    ),
    "DA1": (
        "DA1_研发投入统计表.xlsx",
        [
            ("研发投入汇总",
             ["年份", "营业收入(亿元)", "研发投入(亿元)", "研发强度", "研发人员(人)", "授权专利(件)"],
             [
                 [2020, 198.4, 9.8, "4.9%", 1820, 214],
                 [2021, 218.6, 11.4, "5.2%", 2050, 268],
                 [2022, 231.2, 12.4, "5.4%", 2320, 312],
                 [2023, 242.8, 13.4, "5.5%", 2580, 348],
                 [2024, 250.8, 15.8, "6.3%", 2840, 382],
             ]),
            ("研发项目清单(部分)",
             ["项目编号", "项目名称", "研究方向", "立项年份", "预算(万元)", "进展状态"],
             [
                 ["RD2024-001", "下一代平台核心算法", "核心技术", 2024, 3200, "在研"],
                 ["RD2024-008", "AI赋能产品智能化", "核心技术", 2024, 2800, "在研"],
                 ["RD2023-015", "客户定制解决方案A", "应用技术", 2023, 1200, "已结题"],
                 ["RD2023-022", "量子计算预研", "前沿探索", 2023, 800, "在研"],
                 ["RD2022-031", "新材料轻量化技术", "核心技术", 2022, 2400, "已结题"],
             ]),
        ]
    ),
    "SD1": (
        "SD1_安全生产数据统计.xlsx",
        [
            ("安全绩效指标",
             ["指标名称", "单位", "2022年", "2023年", "2024年", "行业基准(2024)"],
             [
                 ["工伤事故总起数", "起", 8, 6, 4, "-"],
                 ["重大事故起数", "起", 0, 0, 0, "-"],
                 ["TRIR(总记录工伤率)", "‰", 0.68, 0.52, 0.42, "0.85"],
                 ["LTIR(工时损失率)", "‰", 0.24, 0.18, 0.14, "0.35"],
                 ["职业病确诊人数", "人", 0, 0, 0, "-"],
                 ["安全培训覆盖率", "%", 97.2, 98.6, 100, "-"],
                 ["应急演练次数", "次", 22, 28, 30, "-"],
             ]),
            ("事故详情记录",
             ["事故编号", "发生时间", "发生基地", "事故类型", "伤亡情况", "根本原因", "整改措施"],
             [
                 ["ACC2024-01", "2024-02-15", "苏州基地", "轻伤", "1人轻伤", "操作规程未遵守", "加强专项培训"],
                 ["ACC2024-02", "2024-07-08", "上海基地", "轻伤", "1人轻伤", "防护用品未佩戴", "完善检查机制"],
                 ["ACC2024-03", "2024-09-22", "深圳基地", "轻伤", "1人轻伤", "作业环境照明不足", "改善作业环境"],
                 ["ACC2024-04", "2024-11-30", "武汉基地", "轻伤", "1人轻伤", "设备老化未及时维修", "建立设备台账"],
             ]),
        ]
    ),
}


# ── PDF 生成 ──────────────────────────────────────────────────────────────

def make_pdf(path: Path, title: str, sections: list[tuple[str, str]]):
    """生成包含标题+多节正文的 PDF（使用 reportlab，内嵌中文字体）"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 注册系统中文字体
    font_paths = [
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    font_name = "Chinese"
    registered = False
    for fp in font_paths:
        if Path(fp).exists():
            try:
                pdfmetrics.registerFont(TTFont(font_name, fp))
                registered = True
                break
            except Exception:
                continue

    if not registered:
        # 降级：使用 Helvetica，中文字符会显示为方块，但文件合法
        font_name = "Helvetica"

    doc = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=25*mm, rightMargin=25*mm,
        topMargin=20*mm, bottomMargin=20*mm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "title", fontName=font_name, fontSize=16, leading=24,
        alignment=1, spaceAfter=12, textColor=colors.HexColor("#1F4E79"),
    )
    heading_style = ParagraphStyle(
        "heading", fontName=font_name, fontSize=12, leading=18,
        spaceAfter=6, spaceBefore=12, textColor=colors.HexColor("#2E75B6"),
    )
    body_style = ParagraphStyle(
        "body", fontName=font_name, fontSize=10, leading=16,
        spaceAfter=8, firstLineIndent=20,
    )

    story = [
        Paragraph(title, title_style),
        HRFlowable(width="100%", thickness=2, color=colors.HexColor("#1F4E79")),
        Spacer(1, 8*mm),
    ]
    for sec_title, content in sections:
        story.append(Paragraph(sec_title, heading_style))
        story.append(Paragraph(content, body_style))

    doc.build(story)
    print(f"  [PDF ] {path.relative_to(BASE_DIR)}")


PDF_CONTENTS = {
    "GA1": (
        "GA1_公司章程（摘录）.pdf",
        "公司章程（ESG治理相关条款摘录）",
        [
            ("第三章 股东和股东大会",
             "第十八条 公司股东依法享有资产收益、参与重大决策和选择管理者等权利。"
             "股东大会是公司的权力机构，依照法律、行政法规和本章程行使职权。"
             "公司保障每位股东充分行使其合法权益，建立健全股东沟通机制，"
             "通过业绩说明会、投资者开放日等多种方式与股东保持积极互动。"),
            ("第四章 董事会",
             "第三十二条 公司董事会由九名董事组成，设董事长一人。"
             "独立董事不少于董事会成员总数的三分之一。独立董事任期届满前，"
             "非经股东大会审议通过，不得无故解除其职务。"
             "董事会对可持续发展事宜承担最终监督责任，每年至少审议一次ESG战略与绩效。"),
            ("第六章 监事会",
             "第五十八条 公司监事会由五名监事组成，其中职工代表监事不少于三分之一。"
             "监事会对董事、高级管理人员执行公司职务的行为进行监督，对公司财务进行监察。"
             "监事会有权查阅公司账簿，对公司经营情况进行调查。"),
        ]
    ),
    "GC1": (
        "GC1_廉洁合规承诺书（全员签署版）.pdf",
        "廉洁合规承诺书",
        [
            ("承诺内容",
             "本人郑重承诺，在履行岗位职责过程中，严格遵守国家法律法规、公司各项规章制度"
             "及《廉洁从业行为准则》，自觉抵制商业贿赂，不接受、不给予任何形式的不正当利益，"
             "不利用职务便利为本人或他人谋取私利，不泄露公司商业秘密和客户信息。"),
            ("举报机制告知",
             "如发现公司员工存在违规行为，可通过以下渠道进行匿名举报："
             "举报热线：400-XXX-XXXX（全年无休）；举报邮箱：compliance@company.com；"
             "举报信箱：总部大楼一楼大厅合规部信箱。公司承诺对举报人身份严格保密，"
             "不对举报人实施任何形式的打击报复。"),
            ("违规后果",
             "如违反上述承诺，本人愿承担相应的纪律处分，包括但不限于：警告、记过、"
             "降职降薪、解除劳动合同，情节严重者依法移交司法机关处理。"),
        ]
    ),
    "EA3": (
        "EA3_碳减排项目实施报告.pdf",
        "2024年碳减排重点项目实施情况报告",
        [
            ("项目一：工厂屋顶分布式光伏",
             "项目背景：为落实公司碳中和路径规划，2024年在苏州、上海、深圳三个生产基地"
             "推进屋顶分布式光伏扩容工程，新增装机容量12MW，总投资额4,800万元。"
             "实施进展：项目于2024年3月开工，9月全部并网发电，全年发电量1,150万度，"
             "自发自用率88%，预计年均减少碳排放约5,750吨CO₂e，投资回收期约6.8年。"),
            ("项目二：高能耗设备电气化替代",
             "项目背景：对使用天然气、燃油的工业锅炉、加热设备实施电气化改造，"
             "配合绿电采购，实现间接减排。2024年完成改造设备346台套，"
             "涉及压缩空气系统、循环水泵、通风风机等，总投资3,200万元。"
             "减排效果：年节约天然气约320万立方米，折合减少直接排放约6,400吨CO₂e。"),
            ("项目三：绿电采购与绿证",
             "通过广东、江苏电力交易中心完成绿电双边协议交易，2024年采购绿电1.2亿度，"
             "对应绿证120万张，全部注销用于自身碳排放抵消声明。"
             "绿电来源：风电65%（内蒙古、新疆），光伏35%（宁夏、青海）。"),
        ]
    ),
    "DD1": (
        "DD1_网络安全应急预案.pdf",
        "网络安全与数据保护应急预案（2024版）",
        [
            ("一、适用范围与分级",
             "本预案适用于公司所有信息系统、数据资产相关的网络安全事件响应。"
             "安全事件按严重程度分为四级：特别重大（I级）、重大（II级）、较大（III级）、一般（IV级）。"
             "触发I级响应条件：核心系统停服超过4小时，或超过10万条用户数据泄露，"
             "或攻击影响波及业务连续性的关键基础设施。"),
            ("二、响应流程",
             "发现安全事件 → 15分钟内上报信息安全团队 → 启动对应级别响应预案 → "
             "隔离受影响系统 → 取证分析 → 恢复正常运营 → 事后复盘报告。"
             "I/II级事件须在1小时内通知CTO，24小时内向监管部门报告（如适用）。"
             "全程启用应急通信渠道，避免使用可能受影响的企业内网系统。"),
            ("三、恢复目标",
             "RTO（恢复时间目标）：核心业务系统4小时，支撑系统24小时。"
             "RPO（恢复点目标）：核心数据1小时，一般数据24小时。"
             "公司每半年开展一次全规模灾备切换演练，验证RTO/RPO目标可达性。"),
        ]
    ),
    "SA1": (
        "SA1_公益基金会年度报告.pdf",
        "XX公司公益基金会2024年度报告",
        [
            ("基金会概况",
             "XX公司公益基金会成立于2018年，注册资金5,000万元，主要聚焦乡村教育、"
             "科技人才培养和紧急灾害救援三大领域。基金会理事会由9名理事组成，"
             "每年召开两次理事会会议，对重大捐助项目进行审批和监督。"),
            ("2024年度捐助情况",
             "2024年累计捐款3,680万元，其中：乡村教育助学项目1,656万元，"
             "资助贫困学生4,820名，覆盖云南、贵州、四川三省82所学校；"
             "科技人才奖学金项目1,104万元，资助高校优秀理工科学生368名；"
             "灾害救援应急捐款552万元（含新疆洪涝灾害应急捐款）；其他项目368万元。"),
            ("项目成效与评估",
             "乡村教育项目：资助学生平均学业进步率82%，受助学生大学升学率较对照组高21个百分点。"
             "科技奖学金项目：历届获奖者中已有38%进入科研院所或科技企业工作，"
             "人才培养目标持续达成。基金会聘请第三方机构对所有重大项目进行效益评估。"),
        ]
    ),
}


# ── PPTX 生成 ──────────────────────────────────────────────────────────────

def make_pptx(path: Path, title: str, slides: list[tuple[str, list[str]]]):
    """
    生成 PowerPoint 文件。
    slides: [(slide_title, [bullet_text, ...]), ...]
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BLUE  = RGBColor(0x1F, 0x4E, 0x79)
    MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

    blank_layout = prs.slide_layouts[6]  # 完全空白

    def _add_text_box(slide, left, top, width, height, text, font_size=14,
                      bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    # 封面 slide
    cover = prs.slides.add_slide(blank_layout)
    bg = cover.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    _add_text_box(cover, 1, 2.5, 11.33, 1.2, title,
                  font_size=28, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    _add_text_box(cover, 1, 4, 11.33, 0.6, "汇报时间：2025年3月",
                  font_size=14, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)

    # 内容 slides
    for slide_title, bullets in slides:
        sl = prs.slides.add_slide(blank_layout)
        # 顶部色块
        header_bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        header_bg.fill.solid(); header_bg.fill.fore_color.rgb = MID_BLUE
        header_bg.line.fill.background()
        _add_text_box(sl, 0.3, 0.15, 12.5, 0.8, slide_title,
                      font_size=20, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
        # 正文区背景
        body_bg = sl.shapes.add_shape(1, Inches(0.3), Inches(1.3),
                                       Inches(12.73), Inches(5.8))
        body_bg.fill.solid(); body_bg.fill.fore_color.rgb = LIGHT_GRAY
        body_bg.line.fill.background()
        # 子弹点
        for i, bullet in enumerate(bullets):
            _add_text_box(sl, 0.6, 1.5 + i * 1.1, 12.2, 1.0,
                          f"◆  {bullet}", font_size=13,
                          color=RGBColor(0x26,0x26,0x26))

    prs.save(str(path))
    print(f"  [PPTX] {path.relative_to(BASE_DIR)}")


PPTX_CONTENTS = {
    "GA2": (
        "GA2_董事会多元化汇报.pptx",
        "董事会多元化建设专题汇报",
        [
            ("一、多元化现状概览", [
                "董事会规模：9名董事，独立董事3名（占33.3%），符合监管最低要求",
                "性别多元化：女性董事2名（22.2%），较上年提升5.6个百分点",
                "年龄分布：最小42岁，最大67岁，平均52.4岁",
                "国际化背景：具有海外工作或学习经历的董事4名（44.4%）",
            ]),
            ("二、专业背景分布", [
                "财务/审计背景：3名（含2名注册会计师）—— 深化审计委员会专业能力",
                "法律/合规背景：2名（含1名执业律师）—— 强化合规监督功能",
                "行业技术背景：2名（均为高级工程师）—— 提升战略技术判断力",
                "战略/管理背景：2名（均有跨国公司管理经验）—— 增强全球化视野",
            ]),
            ("三、2025-2026提升计划", [
                "目标：到2026年将女性董事比例提升至30%以上（至少新增1名女性独立董事）",
                "行动：委托猎头建立多元化候选人储备库，优先纳入女性及ESG专业背景人选",
                "评估：将多元化目标纳入提名委员会KPI，每年向股东大会披露进展",
                "对标：参考港交所《企业管治守则》C.1.3及MSCI多元化评级标准",
            ]),
        ]
    ),
    "EA1": (
        "EA1_气候战略汇报.pptx",
        "气候变化应对战略与碳中和路径汇报",
        [
            ("一、气候风险识别", [
                "物理风险：极端天气导致基地停产风险（苏州/上海洪涝风险评级：中高）",
                "转型风险：碳税政策收紧（预测2030年碳价：150元/吨，年增成本约500万元）",
                "市场机遇：客户低碳采购要求推动绿色产品溢价（预估3-5%）",
                "监管压力：A股ESG强制披露路线图，2026年起核心指标须经第三方鉴证",
            ]),
            ("二、碳中和承诺与路径", [
                "近期（2024-2026）：范围一+二减排30%（vs 2020基准），绿电占比≥40%",
                "中期（2027-2030）：范围一+二减排50%，完成全部高耗能设备电气化",
                "远期（2031-2050）：全价值链净零排放，启动范围三核算与减排计划",
                "核查机制：每年委托SGS开展第三方碳排放核查，结果随ESG报告公开",
            ]),
            ("三、2024年度执行进展", [
                "范围一+二合计排放30,770 tCO₂e，较基准年下降15.3%，超额完成年度目标",
                "绿电采购1.2亿度，绿电占比35%，较目标差5个百分点（原因：电网绿电紧缺）",
                "光伏装机新增12MW，累计38MW，年发电量3,200万度",
                "碳减排项目实施12项，累计减排约18,400 tCO₂e，投资回报期平均2.2年",
            ]),
        ]
    ),
    "DC1": (
        "DC1_质量管理体系汇报.pptx",
        "产品质量管理体系年度工作汇报",
        [
            ("一、质量管理认证现状", [
                "ISO 9001:2015：覆盖全部8个生产基地，2024年9月完成三年期复评，零不符合项",
                "IATF 16949：覆盖汽车零部件产品线（占营收28%），一次性通过监督审核",
                "客户专项审核：通过大客户A、B、C年度质量体系审核，综合得分≥92分",
                "2025年新目标：推进AS9100D（航空级质量）认证筹备工作",
            ]),
            ("二、质量绩效数据", [
                "产品出厂合格率：99.97%（行业均值：99.82%），连续5年保持99.9%以上",
                "客户退货率：0.018%，较上年下降23%，达到历史最低水平",
                "内部不良品率：1,240 PPM，较上年下降18%",
                "质量损失成本：占销售额0.42%，低于行业均值（0.68%）",
            ]),
            ("三、质量改进重点项目", [
                "供应商质量提升计划：对前50家供应商开展质量能力评级，带动12家完成体系升级",
                "自动化检测升级：引入机器视觉检测系统，覆盖关键工序，误判率下降65%",
                "零缺陷项目：在苏州基地试点，关键产品线实现连续180天零客诉",
                "质量文化建设：全员质量意识培训覆盖率100%，QC小组活动增至86个",
            ]),
        ]
    ),
    "SB1": (
        "SB1_社区沟通汇报.pptx",
        "利益相关方参与与社区沟通年度汇报",
        [
            ("一、利益相关方识别与重要性评估", [
                "2024年开展双向重要性评估，有效问卷2,450份，覆盖6类利益相关方",
                "最高重要性议题：气候变化（94%受访者认为重要）、员工权益（91%）、数据安全（88%）",
                "新兴议题：生物多样性首次进入前10重要议题（排名第7）",
                "评估结果已纳入本年度ESG报告框架与KPI设置",
            ]),
            ("二、社区沟通渠道与效果", [
                "周边居民信访：全年受理128件，100%结案回复，平均响应时间2.8天",
                "月度社区座谈会：参与居民代表平均18人/场，共收集意见建议64条，采纳率61%",
                "年度开放日：苏州基地开放参观，接待周边居民及学生共计840人次",
                "在线互动平台：微信公众号粉丝12.8万，月均留言回复率97%",
            ]),
            ("三、主要议题响应", [
                "噪音投诉：安装隔音屏障，夜间噪音从68dB降至51dB，投诉量下降82%",
                "就业贡献：本地员工占比68%，2024年新招聘本地员工248人",
                "环境污染关切：在线发布月度环境监测数据，透明度大幅提升",
                "2025年计划：建立社区ESG圆桌论坛，每季度开展一次深度对话",
            ]),
        ]
    ),
}


# ── PNG 生成（用 PIL 绘制图表类图片）─────────────────────────────────────

def make_png(path: Path, title: str, chart_type: str, data: dict):
    """
    生成包含简单图表的 PNG 图片（模拟截图/扫描件场景）。
    chart_type: "bar" | "pie" | "text_table"
    """
    from PIL import Image, ImageDraw, ImageFont
    import math

    W, H = 900, 620
    BG    = (255, 255, 255)
    DARK  = (31, 78, 121)
    MID   = (46, 117, 182)
    LIGHT = (189, 215, 238)
    GRAY  = (127, 127, 127)
    BLACK = (30, 30, 30)

    img  = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(img)

    # 顶部色带
    draw.rectangle([0, 0, W, 60], fill=DARK)
    draw.text((W//2, 30), title, fill=(255,255,255), anchor="mm")

    if chart_type == "bar":
        labels  = data["labels"]
        values  = data["values"]
        y_label = data.get("y_label", "")
        bar_w   = min(60, (W - 120) // len(labels))
        max_val = max(values) * 1.1
        chart_top, chart_bot = 100, H - 80
        chart_h = chart_bot - chart_top

        for i, (lbl, val) in enumerate(zip(labels, values)):
            x = 80 + i * (bar_w + 15)
            bar_h = int(chart_h * val / max_val)
            draw.rectangle([x, chart_bot - bar_h, x + bar_w, chart_bot], fill=MID)
            draw.text((x + bar_w//2, chart_bot - bar_h - 14),
                      str(val), fill=BLACK, anchor="mm")
            draw.text((x + bar_w//2, chart_bot + 16), lbl, fill=GRAY, anchor="mt")
        draw.line([60, chart_top, 60, chart_bot], fill=GRAY, width=2)
        draw.line([60, chart_bot, W-40, chart_bot], fill=GRAY, width=2)
        draw.text((30, (chart_top+chart_bot)//2), y_label, fill=GRAY, anchor="mm")

    elif chart_type == "pie":
        labels  = data["labels"]
        values  = data["values"]
        colors_ = [MID, (70,130,180), LIGHT, (173,216,230), (135,206,235),
                   (100,149,237), (176,196,222)]
        total = sum(values)
        cx, cy, r = W//2 - 60, (H+60)//2, 160
        start = -90
        for i, (lbl, val) in enumerate(zip(labels, values)):
            angle = 360 * val / total
            draw.pieslice([cx-r, cy-r, cx+r, cy+r],
                          start=start, end=start+angle, fill=colors_[i % len(colors_)])
            mid_angle = math.radians(start + angle/2)
            tx = int(cx + (r+40) * math.cos(mid_angle))
            ty = int(cy + (r+40) * math.sin(mid_angle))
            draw.text((tx, ty), f"{lbl}\n{val:.1f}%", fill=BLACK, anchor="mm")
            start += angle

    elif chart_type == "text_table":
        headers = data["headers"]
        rows    = data["rows"]
        col_w   = (W - 80) // len(headers)
        draw.rectangle([40, 80, W-40, 112], fill=MID)
        for j, h in enumerate(headers):
            draw.text((40 + j*col_w + col_w//2, 96), h,
                      fill=(255,255,255), anchor="mm")
        for i, row in enumerate(rows):
            y0 = 112 + i * 36
            fill = LIGHT if i % 2 == 0 else BG
            draw.rectangle([40, y0, W-40, y0+36], fill=fill)
            for j, cell in enumerate(row):
                draw.text((40 + j*col_w + col_w//2, y0+18),
                          str(cell), fill=BLACK, anchor="mm")

    # 底部水印
    draw.text((W//2, H-20), "仅供内部测试使用 · 模拟数据", fill=LIGHT, anchor="mm")

    img.save(str(path), "PNG", dpi=(150, 150))
    print(f"  [PNG ] {path.relative_to(BASE_DIR)}")


PNG_CONTENTS = {
    "EA2": (
        "EA2_碳排放趋势图.png", "2020-2024年碳排放趋势（tCO₂e）",
        "bar",
        {"labels": ["2020", "2021", "2022", "2023", "2024"],
         "values": [36300, 35100, 34200, 33000, 30770],
         "y_label": "tCO₂e"},
    ),
    "EC1": (
        "EC1_能源结构饼图.png", "2024年能源消耗结构",
        "pie",
        {"labels": ["电力", "天然气", "柴油", "其他"],
         "values": [62.0, 28.0, 7.0, 3.0]},
    ),
    "SC1": (
        "SC1_员工性别年龄分布图.png", "2024年员工结构分布",
        "bar",
        {"labels": ["<30岁", "30-40岁", "40-50岁", ">50岁"],
         "values": [2703, 4669, 3440, 1474],
         "y_label": "人数"},
    ),
    "DA1": (
        "DA1_研发投入趋势图.png", "2020-2024年研发投入（亿元）",
        "bar",
        {"labels": ["2020", "2021", "2022", "2023", "2024"],
         "values": [9.8, 11.4, 12.4, 13.4, 15.8],
         "y_label": "亿元"},
    ),
    "GC1": (
        "GC1_合规举报处理情况.png", "2024年度合规举报受理与处理情况",
        "text_table",
        {"headers": ["受理类型", "数量", "属实", "不属实", "转办", "处理率"],
         "rows": [
             ["商业贿赂", "3件", "1件", "2件", "0件", "100%"],
             ["利益冲突", "2件", "0件", "1件", "1件", "100%"],
             ["信息泄露", "2件", "1件", "0件", "1件", "100%"],
             ["其他违规", "1件", "0件", "1件", "0件", "100%"],
             ["合计", "8件", "2件", "4件", "2件", "100%"],
         ]},
    ),
    "SD1": (
        "SD1_安全绩效趋势图.png", "2022-2024年安全绩效指标（TRIR）",
        "bar",
        {"labels": ["2022年", "2023年", "2024年", "行业基准"],
         "values": [0.68, 0.52, 0.42, 0.85],
         "y_label": "TRIR(‰)"},
    ),
}


# ── 故意放错位置的文件（模拟甲方操作失误）─────────────────────────────────

MISPLACED_FILES = [
    # 把 EA2 的数据表放到了 EA1（碳排放目标）文件夹 —— 相邻文件夹搞混
    {
        "type": "xlsx",
        "filename": "温室气体排放数据汇总.xlsx",
        "target_folder": FOLDER_MAP["EA1"],  # 错放到 EA1
        "correct_folder": "EA2",
        "sheet_configs": [
            ("排放数据", ["年份", "范围一(tCO₂e)", "范围二(tCO₂e)", "合计"],
             [[2021, 10200, 15800, 26000],
              [2022, 11300, 16900, 28200],
              [2023, 11800, 17500, 29300],
              [2024, 12450, 18320, 30770]]),
            ("分部门排放", ["部门", "范围一(tCO₂e)", "范围二(tCO₂e)"],
             [["生产部", 8200, 12100], ["办公楼", 1800, 4300], ["仓储物流", 2450, 1920]]),
        ],
    },
    # 把 SC2（薪酬福利）的文件放到了 SD1（安全生产）文件夹 —— 编码相近搞混
    {
        "type": "docx",
        "filename": "薪酬市场调研报告.docx",
        "target_folder": FOLDER_MAP["SD1"],  # 错放到 SD1
        "correct_folder": "SC2",
        "title": "2024年薪酬市场调研报告",
        "sections": [
            ("调研背景",
             "本报告基于2024年度行业薪酬调研数据，覆盖制造业、科技行业共1200家样本企业，"
             "旨在评估公司薪酬竞争力水平，为薪酬体系优化提供参考依据。"),
            ("薪酬竞争力分析",
             "公司各岗位薪酬处于行业75分位水平，研发岗位竞争力评分最高（85分位），"
             "生产制造岗位略低（65分位）。与去年相比，整体薪酬竞争力提升5个百分位。"),
            ("薪酬结构",
             "固定薪酬占比68%，绩效奖金占比22%，长期激励占比10%。员工福利涵盖补充医疗保险、"
             "住房公积金（12%）、带薪年假（15天起）及弹性福利积分（每年5000分）。"),
            ("改进建议",
             "建议2025年对生产类岗位薪酬进行专项调整，目标达到行业75分位。"
             "同时优化绩效奖金分配方案，加强与KPI挂钩的透明度，提升员工薪酬满意度。"),
        ],
    },
]


if __name__ == "__main__":
    print("=" * 60)
    print("  模拟甲方文件生成脚本")
    print(f"  目标目录：{BASE_DIR}")
    print("=" * 60)

    if not BASE_DIR.exists():
        print(f"✗ 目标目录不存在：{BASE_DIR}")
        sys.exit(1)

    print("\n[阶段1] 生成 DOCX 文件...")
    for code, folder in FOLDER_MAP.items():
        if code in ALL_DOCX:
            title, sections = ALL_DOCX[code]
            make_docx(folder / f"{code}_{title[:10]}.docx", title, sections)

    print("\n[阶段2a] 生成 XLSX 文件...")
    for code, (filename, sheet_configs) in XLSX_CONTENTS.items():
        make_xlsx(FOLDER_MAP[code] / filename, sheet_configs)

    print("\n[阶段2b] 生成 PDF 文件...")
    for code, (filename, title, sections) in PDF_CONTENTS.items():
        make_pdf(FOLDER_MAP[code] / filename, title, sections)

    print("\n[阶段3a] 生成 PPTX 文件...")
    for code, (filename, title, slides) in PPTX_CONTENTS.items():
        make_pptx(FOLDER_MAP[code] / filename, title, slides)

    print("\n[阶段3b] 生成 PNG 图片...")
    for code, (filename, title, chart_type, chart_data) in PNG_CONTENTS.items():
        make_png(FOLDER_MAP[code] / filename, title, chart_type, chart_data)

    print("\n[阶段4] 放置故意错位的文件（模拟甲方操作失误）...")
    for item in MISPLACED_FILES:
        target = item["target_folder"]
        if item["type"] == "xlsx":
            make_xlsx(target / item["filename"], item["sheet_configs"])
        elif item["type"] == "docx":
            make_docx(target / item["filename"], item["title"], item["sections"])
        print(f"         ↑ 正确位置应为 {item['correct_folder']}，已故意放错")

    print("\n✅ 全部完成！")
    # 统计文件总数
    total = sum(1 for _ in BASE_DIR.rglob("*") if _.is_file() and _.suffix != ".txt")
    print(f"   共生成模拟文件：{total} 个（含 2 个故意错位文件）")
